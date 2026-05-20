"""Build a colourful, infographic-style PPTX in the visual language of
the user's reference template (Copy of Sales Growth Strategy Slides).

Design DNA we carry over from the template:

* Slide size 20 in x 11.25 in (extended widescreen)
* Palette:  #75BBEE  #C7D157  #ECBEFC  #FFB07A  #C6EFFD  on a near-white
            #F8F8F4 background, dark #1A1A1A text
* Fonts:    Radika Next (titles) / Neue Montreal (body) -- with Calibri
            fallback if those are not installed on the presentation
            machine
* Style:    big section numerals in the corner, 3-card layouts with
            numbered (01 / 02 / 03) markers, oversized stat numerals on
            coloured circles, generous whitespace, no footer chrome.

Structure follows the user's brief:
    1.  Title (title only)
    2.  Group members
    3.  Agenda / Table of contents
    ... content slides ...
    last.  Thanks
"""

from __future__ import annotations

import os
from itertools import cycle

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt


# -------- palette  ---------------------------------------------------------

BLUE   = RGBColor(0x75, 0xBB, 0xEE)   # primary accent
OLIVE  = RGBColor(0xC7, 0xD1, 0x57)   # secondary
LAVEND = RGBColor(0xEC, 0xBE, 0xFC)   # tertiary
PEACH  = RGBColor(0xFF, 0xB0, 0x7A)   # highlight
SKY    = RGBColor(0xC6, 0xEF, 0xFD)   # background tint
INK    = RGBColor(0x1A, 0x1A, 0x1A)   # body text
GRAY   = RGBColor(0x55, 0x55, 0x55)   # secondary text
PAGE   = RGBColor(0xF8, 0xF8, 0xF4)   # page background
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DEEP   = RGBColor(0x1B, 0x2C, 0x4A)   # deep navy for hero titles

ACCENTS = [BLUE, OLIVE, LAVEND, PEACH, SKY]


# -------- geometry ---------------------------------------------------------

SLIDE_W = Inches(20)
SLIDE_H = Inches(11.25)
MARGIN  = Inches(0.9)
SECTION_NUMERAL_LEFT = Inches(0.9)
SECTION_NUMERAL_TOP  = Inches(0.7)


FIG_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                       "report", "figs")
OUT_PATH = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                        "AO_PSO_Presentation_v2.pptx")


def fig(name: str) -> str:
    return os.path.join(FIG_DIR, name)


# -------- font helpers -----------------------------------------------------

TITLE_FONT = "Radika Next"
BODY_FONT  = "Neue Montreal"


def set_run(run, *, size=None, bold=None, color=None, italic=None, font=None):
    if font is not None:
        run.font.name = font
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if italic is not None:
        run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color


def shape_no_outline(sh):
    sh.line.fill.background()
    sh.shadow.inherit = False


def fill_solid(sh, color):
    sh.fill.solid()
    sh.fill.fore_color.rgb = color


def add_text(slide, text, *, left, top, width, height,
             size, bold=False, italic=False, color=INK,
             font=BODY_FONT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    set_run(r, size=size, bold=bold, italic=italic, color=color, font=font)
    return tb


def add_multiline(slide, lines, *, left, top, width, height,
                   anchor=MSO_ANCHOR.TOP):
    """``lines`` is a list of ``(text, opts)`` tuples.

    Recognised opts:  size, bold, italic, color, font, align, indent,
                       space_after.
    """
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    first = True
    for entry in lines:
        if isinstance(entry, tuple):
            text, opts = entry
        else:
            text, opts = entry, {}
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = opts.get("align", PP_ALIGN.LEFT)
        p.level = opts.get("indent", 0)
        p.space_after = Pt(opts.get("space_after", 6))
        r = p.add_run()
        r.text = text
        set_run(r,
                size=opts.get("size", 20),
                bold=opts.get("bold", False),
                italic=opts.get("italic", False),
                color=opts.get("color", INK),
                font=opts.get("font", BODY_FONT))
    return tb


def add_page_background(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    shape_no_outline(bg)
    fill_solid(bg, PAGE)
    # send to back: python-pptx puts shapes in z-order added; we add bg first
    return bg


def add_section_numeral(slide, numeral: str, color=BLUE):
    """Big light section numeral in the corner -- the template's
    signature visual hook."""
    tb = slide.shapes.add_textbox(SECTION_NUMERAL_LEFT,
                                   SECTION_NUMERAL_TOP,
                                   Inches(2.2), Inches(2.0))
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_top = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = numeral
    set_run(r, size=120, bold=True, color=color, font=TITLE_FONT)


def add_accent_dot(slide, left, top, diameter, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diameter, diameter)
    shape_no_outline(sh)
    fill_solid(sh, color)
    return sh


def add_corner_decor(slide, accent=BLUE):
    """Decorative dots / circles top-right corner of every content slide."""
    # one big translucent circle clipped at the corner
    big = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 SLIDE_W - Inches(2.8),
                                 -Inches(1.6),
                                 Inches(4.5), Inches(4.5))
    shape_no_outline(big)
    fill_solid(big, accent)
    # ring of small dots
    for i, c in enumerate([OLIVE, LAVEND, PEACH]):
        add_accent_dot(slide, SLIDE_W - Inches(0.6) - Inches(0.4)*i,
                       SLIDE_H - Inches(0.7),
                       Inches(0.28), c)


# -------- slide builders ---------------------------------------------------

def slide_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_page_background(s)

    # Big decorative circles in the corners
    c1 = s.shapes.add_shape(MSO_SHAPE.OVAL,
                            -Inches(3.5), -Inches(3.5),
                            Inches(9.0), Inches(9.0))
    shape_no_outline(c1); fill_solid(c1, BLUE)

    c2 = s.shapes.add_shape(MSO_SHAPE.OVAL,
                            SLIDE_W - Inches(5.0),
                            SLIDE_H - Inches(5.0),
                            Inches(7.0), Inches(7.0))
    shape_no_outline(c2); fill_solid(c2, OLIVE)

    c3 = s.shapes.add_shape(MSO_SHAPE.OVAL,
                            SLIDE_W - Inches(8.5),
                            -Inches(1.5),
                            Inches(2.5), Inches(2.5))
    shape_no_outline(c3); fill_solid(c3, LAVEND)

    c4 = s.shapes.add_shape(MSO_SHAPE.OVAL,
                            Inches(1.5),
                            SLIDE_H - Inches(3.0),
                            Inches(1.8), Inches(1.8))
    shape_no_outline(c4); fill_solid(c4, PEACH)

    # small accent dots
    for x, y, c in [
        (Inches(3.3), Inches(2.2), PEACH),
        (Inches(17.0), Inches(3.5), SKY),
        (Inches(15.5), Inches(8.5), LAVEND),
    ]:
        add_accent_dot(s, x, y, Inches(0.35), c)

    # central title block on white card
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(2.8), Inches(2.6),
                              Inches(14.4), Inches(6.0))
    card.adjustments[0] = 0.04
    shape_no_outline(card); fill_solid(card, WHITE)

    add_text(s, "FINAL PHASE PROJECT", left=Inches(3.6), top=Inches(3.1),
             width=Inches(12), height=Inches(0.5),
             size=20, bold=True, color=BLUE, font=BODY_FONT)

    add_text(s, "AO-PSO", left=Inches(3.6), top=Inches(3.6),
             width=Inches(13), height=Inches(2.0),
             size=110, bold=True, color=DEEP, font=TITLE_FONT)

    add_text(s,
             "Adaptive Opposition-Based Particle Swarm Optimization",
             left=Inches(3.6), top=Inches(5.6),
             width=Inches(13), height=Inches(0.9),
             size=30, bold=False, color=INK, font=TITLE_FONT)

    add_text(s,
             "with Chaotic Initialization, Generalized Opposites,\n"
             "and Diversity-Driven Generation Jumping",
             left=Inches(3.6), top=Inches(6.4),
             width=Inches(13), height=Inches(1.4),
             size=20, italic=True, color=GRAY, font=BODY_FONT)

    add_text(s,
             "Swarm Intelligence  /  National University of Technology",
             left=Inches(3.6), top=Inches(7.85),
             width=Inches(13), height=Inches(0.4),
             size=14, bold=True, color=GRAY, font=BODY_FONT)

    s.notes_slide.notes_text_frame.text = (
        "Aleena opens. Title only -- introduce the project name; do not "
        "list members yet (next slide).")


# ---------- group members slide --------------------------------------------

def slide_members(prs):
    members = [
        ("Aleena Tahir",   "Roll 1",  "IEEE compilation"),
        ("Eman Asghar",    "Roll 10", "Diagrams & visuals"),
        ("Aena Habib",     "Roll 20", "Value addition lead"),
        ("Malaika Akhter", "Roll 22", "Abstract & intro"),
        ("Dua Kamal",      "Roll 23", "Presentation design"),
        ("Sadia Mazhar",   "Roll 29", "Related work"),
        ("Inshal Adil",    "Roll 38", "Results analysis"),
        ("Saqlain Abbas",  "Roll 48", "Implementation"),
    ]
    palette = [BLUE, OLIVE, LAVEND, PEACH, SKY, BLUE, OLIVE, LAVEND]

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_background(s)

    add_section_numeral(s, "01", color=BLUE)
    add_text(s, "MEET THE TEAM",
             left=Inches(3.8), top=Inches(1.0),
             width=Inches(14), height=Inches(0.5),
             size=18, bold=True, color=BLUE, font=BODY_FONT)
    add_text(s, "Group members",
             left=Inches(3.8), top=Inches(1.45),
             width=Inches(14), height=Inches(1.2),
             size=64, bold=True, color=DEEP, font=TITLE_FONT)
    add_text(s,
             "Department of Artificial Intelligence  /  Semester 6  /  "
             "Instructor: Lec. Faria Sajjad",
             left=Inches(3.8), top=Inches(2.65),
             width=Inches(14), height=Inches(0.5),
             size=16, italic=True, color=GRAY, font=BODY_FONT)

    # 4 x 2 grid of member cards
    card_w = Inches(4.2)
    card_h = Inches(3.0)
    gap_x = Inches(0.35)
    gap_y = Inches(0.4)
    grid_left = (SLIDE_W - (4 * card_w + 3 * gap_x)) // 2
    grid_top = Inches(3.8)

    for idx, (name, roll, role) in enumerate(members):
        row = idx // 4
        col = idx % 4
        left = grid_left + col * (card_w + gap_x)
        top = grid_top + row * (card_h + gap_y)

        bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                left, top, card_w, card_h)
        bg.adjustments[0] = 0.08
        shape_no_outline(bg)
        fill_solid(bg, WHITE)

        # colored top stripe
        stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    left, top, card_w, Inches(0.7))
        shape_no_outline(stripe)
        fill_solid(stripe, palette[idx])

        # initial circle
        initial = name.split()[0][0] + name.split()[-1][0]
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                  left + Inches(0.3),
                                  top + Inches(1.0),
                                  Inches(1.3), Inches(1.3))
        shape_no_outline(dot)
        fill_solid(dot, palette[idx])
        # initials inside
        ini_tb = s.shapes.add_textbox(left + Inches(0.3),
                                       top + Inches(1.0),
                                       Inches(1.3), Inches(1.3))
        ini_tf = ini_tb.text_frame
        ini_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ini_tf.margin_left = 0; ini_tf.margin_right = 0
        ini_tf.margin_top = 0; ini_tf.margin_bottom = 0
        pp = ini_tf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
        rr = pp.add_run(); rr.text = initial
        set_run(rr, size=30, bold=True, color=WHITE, font=TITLE_FONT)

        # roll badge top-right (pill on light background so it stays readable)
        add_text(s, roll,
                 left=left + Inches(1.8), top=top + Inches(1.05),
                 width=card_w - Inches(1.9), height=Inches(0.4),
                 size=14, bold=True, color=GRAY, font=BODY_FONT)

        # name + role
        add_text(s, name,
                 left=left + Inches(0.3), top=top + Inches(2.35),
                 width=card_w - Inches(0.6), height=Inches(0.5),
                 size=20, bold=True, color=DEEP, font=TITLE_FONT)
        add_text(s, role,
                 left=left + Inches(0.3), top=top + Inches(2.65),
                 width=card_w - Inches(0.6), height=Inches(0.35),
                 size=12, color=GRAY, italic=True, font=BODY_FONT)

    s.notes_slide.notes_text_frame.text = (
        "Aleena -- one quick sentence per member's responsibility.")


# ---------- agenda ---------------------------------------------------------

def slide_agenda(prs):
    items = [
        ("01", "Problem & motivation",          BLUE),
        ("02", "Related work & focal paper",    OLIVE),
        ("03", "Our seven value additions",     LAVEND),
        ("04", "Architecture & flow",           PEACH),
        ("05", "Experimental setup",            SKY),
        ("06", "Results & significance",        BLUE),
        ("07", "Honest limitation: Schwefel",   OLIVE),
        ("08", "Conclusion & Q&A",              LAVEND),
    ]

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_background(s)
    add_section_numeral(s, "02", color=OLIVE)

    add_text(s, "WHAT YOU WILL SEE",
             left=Inches(3.8), top=Inches(1.0),
             width=Inches(14), height=Inches(0.5),
             size=18, bold=True, color=OLIVE, font=BODY_FONT)
    add_text(s, "Agenda",
             left=Inches(3.8), top=Inches(1.45),
             width=Inches(14), height=Inches(1.2),
             size=64, bold=True, color=DEEP, font=TITLE_FONT)

    # two-column list of pill-style entries
    col1 = items[:4]
    col2 = items[4:]
    pill_w = Inches(7.4)
    pill_h = Inches(1.15)
    gap = Inches(0.3)
    col1_left = Inches(2.2)
    col2_left = SLIDE_W - col1_left - pill_w
    col_top = Inches(3.7)

    for col, left in [(col1, col1_left), (col2, col2_left)]:
        for j, (num, label, color) in enumerate(col):
            top = col_top + j * (pill_h + gap)
            pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       left, top, pill_w, pill_h)
            pill.adjustments[0] = 0.45
            shape_no_outline(pill)
            fill_solid(pill, WHITE)

            # left numbered circle
            dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                      left + Inches(0.18),
                                      top + Inches(0.18),
                                      Inches(0.8), Inches(0.8))
            shape_no_outline(dot)
            fill_solid(dot, color)
            n_tb = s.shapes.add_textbox(left + Inches(0.18),
                                        top + Inches(0.18),
                                        Inches(0.8), Inches(0.8))
            n_tf = n_tb.text_frame
            n_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            n_tf.margin_left = 0; n_tf.margin_right = 0
            n_tf.margin_top = 0; n_tf.margin_bottom = 0
            pp = n_tf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
            rr = pp.add_run(); rr.text = num
            set_run(rr, size=18, bold=True, color=WHITE, font=TITLE_FONT)

            add_text(s, label,
                     left=left + Inches(1.2), top=top + Inches(0.30),
                     width=pill_w - Inches(1.3), height=Inches(0.6),
                     size=22, bold=True, color=DEEP, font=TITLE_FONT)

    s.notes_slide.notes_text_frame.text = (
        "Aleena -- 30 seconds; tell the audience what to expect.")


# ---------- generic content scaffolding ------------------------------------

def add_content_header(slide, section_num, section_label, big_title,
                        section_color=BLUE):
    add_section_numeral(slide, section_num, color=section_color)
    add_text(slide, section_label,
             left=Inches(3.8), top=Inches(1.0),
             width=Inches(15), height=Inches(0.5),
             size=18, bold=True, color=section_color, font=BODY_FONT)
    add_text(slide, big_title,
             left=Inches(3.8), top=Inches(1.45),
             width=Inches(15), height=Inches(1.4),
             size=54, bold=True, color=DEEP, font=TITLE_FONT)


def add_subtitle(slide, text, top=Inches(2.85)):
    add_text(slide, text,
             left=Inches(3.8), top=top,
             width=Inches(15.5), height=Inches(0.5),
             size=18, italic=True, color=GRAY, font=BODY_FONT)


# ---------- slide: problem and motivation ----------------------------------

def slide_problem(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_background(s)
    add_content_header(s, "01", "INTRODUCTION",
                        "Why we worked on opposition-based PSO",
                        section_color=BLUE)

    # two-card layout: failure modes + the OBL idea
    cards = [
        (BLUE,  "TWO FAILURE MODES OF PSO",
         "•  Premature convergence -- swarm collapses\n"
         "    into a sub-optimal basin\n\n"
         "•  Sensitivity to initialization -- pseudo-random\n"
         "    seeding gives uneven coverage of the\n"
         "    search space"),
        (OLIVE, "OPPOSITION-BASED LEARNING (Tizhoosh, 2005)",
         "Whenever you evaluate x,\n"
         "also evaluate its opposite  x̂ = a + b − x.\n\n"
         "One of the two is, on average, closer to\n"
         "any unknown optimum -- so OBL halves the\n"
         "expected distance to it."),
    ]
    add_two_cards(s, cards, top=Inches(3.8))

    add_text(s,
             "Goal: push OBL deeper into the algorithm and report it honestly.",
             left=Inches(2.0), top=Inches(9.8),
             width=Inches(16), height=Inches(0.6),
             size=26, bold=True, color=DEEP, font=TITLE_FONT,
             align=PP_ALIGN.CENTER)

    s.notes_slide.notes_text_frame.text = (
        "Malaika -- explain motivation; emphasise the two failure modes.")


def add_two_cards(s, cards, *, top, height=Inches(5.8)):
    n = len(cards)
    gap = Inches(0.6)
    total_w = SLIDE_W - 2 * Inches(2.0)
    card_w = (total_w - (n - 1) * gap) // n
    left0 = Inches(2.0)
    for i, (color, title, body) in enumerate(cards):
        left = left0 + i * (card_w + gap)
        bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 left, top, card_w, height)
        bg.adjustments[0] = 0.05
        shape_no_outline(bg)
        fill_solid(bg, WHITE)
        stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     left, top, card_w, Inches(0.45))
        shape_no_outline(stripe)
        fill_solid(stripe, color)
        add_text(s, title,
                 left=left + Inches(0.5), top=top + Inches(0.75),
                 width=card_w - Inches(1.0), height=Inches(0.8),
                 size=16, bold=True, color=color, font=BODY_FONT)
        add_text(s, body,
                 left=left + Inches(0.5), top=top + Inches(1.55),
                 width=card_w - Inches(1.0), height=height - Inches(1.7),
                 size=20, color=INK, font=BODY_FONT)


def add_three_cards(s, cards, *, top, height=Inches(5.5)):
    n = 3
    gap = Inches(0.4)
    total_w = SLIDE_W - 2 * Inches(1.4)
    card_w = (total_w - (n - 1) * gap) // n
    left0 = Inches(1.4)
    for i, (num, color, title, body) in enumerate(cards):
        left = left0 + i * (card_w + gap)

        bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 left, top, card_w, height)
        bg.adjustments[0] = 0.05
        shape_no_outline(bg); fill_solid(bg, WHITE)

        # numbered chip
        chip = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                   left + Inches(0.5),
                                   top + Inches(0.5),
                                   Inches(1.2), Inches(1.2))
        shape_no_outline(chip); fill_solid(chip, color)
        chip_tb = s.shapes.add_textbox(left + Inches(0.5),
                                        top + Inches(0.5),
                                        Inches(1.2), Inches(1.2))
        chip_tf = chip_tb.text_frame
        chip_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        chip_tf.margin_left = 0; chip_tf.margin_right = 0
        chip_tf.margin_top = 0; chip_tf.margin_bottom = 0
        pp = chip_tf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
        rr = pp.add_run(); rr.text = num
        set_run(rr, size=26, bold=True, color=WHITE, font=TITLE_FONT)

        add_text(s, title,
                 left=left + Inches(0.5), top=top + Inches(2.0),
                 width=card_w - Inches(1.0), height=Inches(0.8),
                 size=24, bold=True, color=DEEP, font=TITLE_FONT)
        add_text(s, body,
                 left=left + Inches(0.5), top=top + Inches(2.85),
                 width=card_w - Inches(1.0), height=height - Inches(3.0),
                 size=17, color=GRAY, font=BODY_FONT)


# ---------- related work ---------------------------------------------------

def slide_related(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_background(s)
    add_content_header(s, "02", "RELATED WORK",
                        "Fifteen successor papers in the OBL-PSO line",
                        section_color=OLIVE)

    add_subtitle(s, "We grouped them into four families to keep the comparison readable.")

    families = [
        ("01", BLUE,
         "OBL at initialization",
         "O-PSO 2009 (focal),\nCSPSO, WE-PSO,\nLatinPSO, Sobol-PSO"),
        ("02", OLIVE,
         "OBL throughout the search",
         "OPSO-Cauchy 2007,\nOVCPSO 2009, GOPSO 2011,\nBPSO-OBL 2015"),
        ("03", LAVEND,
         "PSO variants with OBL",
         "OCLPSO, IOBL-PSO,\nCOPSO, PSO-OBL-NNI,\nMutation-OBL-PSO"),
        ("04", PEACH,
         "Beyond PSO / applied",
         "ODE (DE), SDE-OBL,\nAMO-OBL,\nReactive-power dispatch"),
    ]
    add_four_cards(s, families, top=Inches(3.6), height=Inches(4.6))

    # the 4 open problems banner at bottom
    add_banner(s,
               "Open problems flagged across these papers:",
               ("Fixed Jr  •  Linear opposite only  •  No FE-fair accounting  •  No statistical tests"),
               top=Inches(8.6), color=DEEP)

    s.notes_slide.notes_text_frame.text = (
        "Sadia -- read the four families and the four open problems below. "
        "Those open problems map one-to-one onto our value additions.")


def add_four_cards(s, cards, *, top, height=Inches(4.5)):
    n = 4
    gap = Inches(0.3)
    total_w = SLIDE_W - 2 * Inches(1.4)
    card_w = (total_w - (n - 1) * gap) // n
    left0 = Inches(1.4)
    for i, (num, color, title, body) in enumerate(cards):
        left = left0 + i * (card_w + gap)
        bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 left, top, card_w, height)
        bg.adjustments[0] = 0.05
        shape_no_outline(bg); fill_solid(bg, WHITE)
        chip = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                   left + Inches(0.4),
                                   top + Inches(0.45),
                                   Inches(1.0), Inches(1.0))
        shape_no_outline(chip); fill_solid(chip, color)
        chip_tb = s.shapes.add_textbox(left + Inches(0.4),
                                        top + Inches(0.45),
                                        Inches(1.0), Inches(1.0))
        chip_tf = chip_tb.text_frame
        chip_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        chip_tf.margin_left = 0; chip_tf.margin_right = 0
        chip_tf.margin_top = 0; chip_tf.margin_bottom = 0
        pp = chip_tf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
        rr = pp.add_run(); rr.text = num
        set_run(rr, size=22, bold=True, color=WHITE, font=TITLE_FONT)
        add_text(s, title,
                 left=left + Inches(0.4), top=top + Inches(1.7),
                 width=card_w - Inches(0.8), height=Inches(0.7),
                 size=20, bold=True, color=DEEP, font=TITLE_FONT)
        add_text(s, body,
                 left=left + Inches(0.4), top=top + Inches(2.45),
                 width=card_w - Inches(0.8), height=height - Inches(2.5),
                 size=14, color=GRAY, font=BODY_FONT)


def add_banner(s, title, body, *, top, color=BLUE):
    bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(1.4), top,
                             SLIDE_W - Inches(2.8), Inches(1.8))
    bg.adjustments[0] = 0.2
    shape_no_outline(bg); fill_solid(bg, color)
    add_text(s, title,
             left=Inches(2.0), top=top + Inches(0.25),
             width=SLIDE_W - Inches(4.0), height=Inches(0.5),
             size=18, bold=True, color=WHITE, font=BODY_FONT,
             align=PP_ALIGN.CENTER)
    add_text(s, body,
             left=Inches(2.0), top=top + Inches(0.85),
             width=SLIDE_W - Inches(4.0), height=Inches(0.8),
             size=22, bold=True, color=WHITE, font=TITLE_FONT,
             align=PP_ALIGN.CENTER)


# ---------- focal paper ----------------------------------------------------

def slide_focal(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_background(s)
    add_content_header(s, "02", "THE FOCAL PAPER",
                        "O-PSO — Jabeen, Jalil & Baig, GECCO 2009",
                        section_color=OLIVE)
    add_subtitle(s, "What it did, and where it stops")

    # left card: what it does (numbered steps)
    left_top = Inches(3.6)
    card_h = Inches(6.5)
    card_w = Inches(7.6)
    left_left = Inches(1.4)

    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               left_left, left_top, card_w, card_h)
    card.adjustments[0] = 0.04
    shape_no_outline(card); fill_solid(card, WHITE)
    stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 left_left, left_top, Inches(0.25), card_h)
    shape_no_outline(stripe); fill_solid(stripe, BLUE)

    add_text(s, "What O-PSO does",
             left=left_left + Inches(0.6), top=left_top + Inches(0.4),
             width=card_w - Inches(0.8), height=Inches(0.7),
             size=28, bold=True, color=DEEP, font=TITLE_FONT)
    steps = [
        ("1.", "Generate N random particles"),
        ("2.", "Compute N linear opposites   x̂ = a + b − x"),
        ("3.", "Evaluate all 2N candidates; keep the best N"),
        ("4.", "Run standard Clerc-constriction PSO\n      until the budget is exhausted"),
    ]
    for i, (n, t) in enumerate(steps):
        y = left_top + Inches(1.5) + Inches(1.05) * i
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                   left_left + Inches(0.6), y,
                                   Inches(0.6), Inches(0.6))
        shape_no_outline(circ); fill_solid(circ, BLUE)
        n_tb = s.shapes.add_textbox(left_left + Inches(0.6), y,
                                    Inches(0.6), Inches(0.6))
        n_tb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        n_tb.text_frame.margin_left = 0; n_tb.text_frame.margin_right = 0
        n_tb.text_frame.margin_top = 0; n_tb.text_frame.margin_bottom = 0
        pp = n_tb.text_frame.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
        rr = pp.add_run(); rr.text = n.replace(".", "")
        set_run(rr, size=18, bold=True, color=WHITE, font=TITLE_FONT)
        add_text(s, t,
                 left=left_left + Inches(1.4), top=y + Inches(0.06),
                 width=card_w - Inches(1.6), height=Inches(0.9),
                 size=17, color=INK, font=BODY_FONT)

    # right card: limitations
    right_left = left_left + card_w + Inches(0.4)
    card2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                right_left, left_top,
                                SLIDE_W - right_left - Inches(1.4), card_h)
    card2.adjustments[0] = 0.04
    shape_no_outline(card2); fill_solid(card2, WHITE)
    stripe2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  right_left, left_top, Inches(0.25), card_h)
    shape_no_outline(stripe2); fill_solid(stripe2, PEACH)

    add_text(s, "Limitations we will address",
             left=right_left + Inches(0.6), top=left_top + Inches(0.4),
             width=Inches(8.0), height=Inches(0.7),
             size=28, bold=True, color=DEEP, font=TITLE_FONT)

    cons = [
        ("•",  "Pseudo-random seed -- no guaranteed space coverage"),
        ("•",  "Linear opposite only -- geometrically rigid"),
        ("•",  "No OBL after generation 0 -- most of the search runs blind"),
        ("•",  "No Wilcoxon, no FE-budget reporting -- claims hard to verify"),
    ]
    for i, (m, t) in enumerate(cons):
        y = left_top + Inches(1.5) + Inches(1.05) * i
        bullet = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                     right_left + Inches(0.7), y + Inches(0.15),
                                     Inches(0.3), Inches(0.3))
        shape_no_outline(bullet); fill_solid(bullet, PEACH)
        add_text(s, t,
                 left=right_left + Inches(1.2), top=y + Inches(0.08),
                 width=Inches(8.5), height=Inches(0.85),
                 size=17, color=INK, font=BODY_FONT)

    s.notes_slide.notes_text_frame.text = (
        "Aena -- O-PSO is clean and minimal. Everything after this slide is "
        "our improvement on top of it.")


# ---------- value additions overview ---------------------------------------

def slide_va_overview(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_background(s)
    add_content_header(s, "03", "OUR VALUE ADDITIONS",
                        "Seven enhancements, one novel mechanism",
                        section_color=LAVEND)
    add_subtitle(s, "Six borrowed from the literature, one new -- the adaptive jumping rate.")

    # 7 small infographic chips arranged in two rows
    rows = [
        # row 1 (4 chips)
        [("VA-1", BLUE,  "Chaotic init",          "logistic map"),
         ("VA-2", OLIVE, "Generalized OBL",       "GOBL formula"),
         ("VA-3", LAVEND, "Generation jumping",   "every iteration"),
         ("VA-4", PEACH, "Adaptive Jr",           "★ NOVEL")],
        # row 2 (3 chips centred)
        [("VA-5", SKY,   "pbest opposition",      "on stagnation"),
         ("VA-6", BLUE,  "Rebel repulsion",       "away from worst"),
         ("VA-7", OLIVE, "FE + Wilcoxon",         "honest reporting")],
    ]

    chip_w = Inches(3.9)
    chip_h = Inches(2.7)
    gap = Inches(0.3)

    for row_i, row in enumerate(rows):
        n = len(row)
        total_w = n * chip_w + (n - 1) * gap
        row_left = (SLIDE_W - total_w) // 2
        row_top = Inches(3.8) + Inches(0.0) + row_i * (chip_h + Inches(0.4))
        for j, (num, color, title, sub) in enumerate(row):
            left = row_left + j * (chip_w + gap)
            top = row_top
            bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     left, top, chip_w, chip_h)
            bg.adjustments[0] = 0.08
            shape_no_outline(bg)
            novel = "NOVEL" in sub
            fill_solid(bg, color if not novel else PEACH)
            # white card overlay if not novel
            if not novel:
                inner = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                            left + Inches(0.08),
                                            top + Inches(0.08),
                                            chip_w - Inches(0.16),
                                            chip_h - Inches(0.16))
                inner.adjustments[0] = 0.08
                shape_no_outline(inner); fill_solid(inner, WHITE)

            text_color = WHITE if novel else color
            add_text(s, num,
                     left=left + Inches(0.4), top=top + Inches(0.3),
                     width=chip_w - Inches(0.8), height=Inches(0.5),
                     size=22, bold=True, color=text_color, font=BODY_FONT)
            title_color = WHITE if novel else DEEP
            add_text(s, title,
                     left=left + Inches(0.4), top=top + Inches(0.95),
                     width=chip_w - Inches(0.8), height=Inches(0.9),
                     size=26, bold=True, color=title_color, font=TITLE_FONT)
            sub_color = WHITE if novel else GRAY
            add_text(s, sub,
                     left=left + Inches(0.4), top=top + Inches(1.9),
                     width=chip_w - Inches(0.8), height=Inches(0.6),
                     size=16, italic=True, color=sub_color, font=BODY_FONT)

    s.notes_slide.notes_text_frame.text = (
        "Aena -- the peach (VA-4) chip is the novel one. The other six are "
        "principled re-uses from the literature, packaged together for the "
        "first time inside the O-PSO plug-in framework.")


# ---------- VA-1 + VA-2 ----------------------------------------------------

def slide_va_init(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_background(s)
    add_content_header(s, "03", "VA-1 + VA-2",
                        "Better starting points",
                        section_color=BLUE)
    add_subtitle(s, "Replace the random seed and the linear opposite at generation zero.")

    cards = [
        (BLUE,
         "VA-1   ·   Chaotic logistic-map init",
         "x_{n+1} = μ · x_n · (1 − x_n),  μ = 4\n\n"
         "One independent chaotic stream per\n"
         "dimension. Ergodic distribution covers\n"
         "the unit interval more uniformly than the\n"
         "Mersenne Twister at small sample sizes."),
        (OLIVE,
         "VA-2   ·   Generalized Opposite (GOBL)",
         "x̂_j  =  a_j + b_j − β · x_j − (1 − β) · x̄_j\n\n"
         "β ~ U(0,1) per particle; x̄ is the\n"
         "per-dimension swarm mean. The mean-aware\n"
         "term breaks box symmetry, so opposites\n"
         "are no longer simple reflections."),
    ]
    add_two_cards(s, cards, top=Inches(3.6), height=Inches(6.3))

    add_text(s,
             "Cost stays at 2N initial evaluations -- exactly the same budget as O-PSO.",
             left=Inches(2.0), top=Inches(10.2),
             width=Inches(16), height=Inches(0.5),
             size=20, bold=True, color=OLIVE, font=BODY_FONT,
             align=PP_ALIGN.CENTER, italic=True)

    s.notes_slide.notes_text_frame.text = (
        "Saqlain -- chaotic init costs nothing extra; GOBL is a one-line "
        "formula change but a much richer opposite.")


# ---------- VA-3 -----------------------------------------------------------

def slide_va_gen_jumping(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_background(s)
    add_content_header(s, "03", "VA-3",
                        "Generation jumping",
                        section_color=BLUE)
    add_subtitle(s, "Apply OBL throughout the search, not just at initialisation.")

    # left: explanation; right: the algorithm steps card
    add_multiline(s, [
        ("In O-PSO, OBL fires only once -- at generation 0.",
         {"size": 22, "bold": True, "color": DEEP, "font": TITLE_FONT,
          "space_after": 12}),
        ("After that the swarm has only random init to thank for any diversity.",
         {"size": 19, "color": GRAY, "italic": True, "space_after": 24}),
        ("Generation jumping fixes this:",
         {"size": 22, "bold": True, "color": BLUE, "font": TITLE_FONT,
          "space_after": 12}),
        ("With probability Jr at every generation,",
         {"size": 19, "color": INK}),
        ("    1.  compute  X̂ = GOBL(X)",
         {"size": 19, "color": INK}),
        ("    2.  evaluate all N opposites",
         {"size": 19, "color": INK}),
        ("    3.  replace x_i with x̂_i wherever f(x̂_i) < f(x_i).",
         {"size": 19, "color": INK, "space_after": 24}),
        ("Cost:  +N FEs each time the jump fires",
         {"size": 18, "italic": True, "color": GRAY,
          "space_after": 24}),
        ("Open problem in the literature:  Jr is always picked by hand.",
         {"size": 22, "bold": True, "color": PEACH, "font": TITLE_FONT}),
        ("We fix that next  ►",
         {"size": 22, "bold": True, "color": PEACH, "font": TITLE_FONT}),
    ], left=Inches(2.0), top=Inches(3.7),
        width=Inches(9.5), height=Inches(7.0))

    # right: visual chip
    chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(12.5), Inches(3.7),
                               Inches(6.5), Inches(6.8))
    chip.adjustments[0] = 0.05
    shape_no_outline(chip); fill_solid(chip, LAVEND)
    add_text(s, "Probability of jump",
             left=Inches(12.9), top=Inches(4.1),
             width=Inches(6.0), height=Inches(0.6),
             size=18, bold=True, color=DEEP, font=BODY_FONT,
             align=PP_ALIGN.CENTER)
    add_text(s, "Jr",
             left=Inches(12.9), top=Inches(4.9),
             width=Inches(6.0), height=Inches(1.4),
             size=120, bold=True, color=DEEP, font=TITLE_FONT,
             align=PP_ALIGN.CENTER)
    add_text(s,
             "Same value every generation.\n"
             "Sometimes too low,\n"
             "sometimes wasteful.",
             left=Inches(12.9), top=Inches(7.6),
             width=Inches(6.0), height=Inches(2.5),
             size=18, color=GRAY, italic=True, font=BODY_FONT,
             align=PP_ALIGN.CENTER)

    s.notes_slide.notes_text_frame.text = (
        "Saqlain -- explain that without an adaptive Jr the FE budget is "
        "wasted on OBL late in the search when the swarm has already "
        "converged.")


# ---------- VA-4 NOVEL ----------------------------------------------------

def slide_va_novel(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_background(s)
    add_content_header(s, "03", "VA-4   ★   NOVEL",
                        "Adaptive jumping rate",
                        section_color=PEACH)
    add_subtitle(s, "Tie Jr to a live measurement of swarm diversity.")

    # the big formula in a peach band
    band = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(1.4), Inches(3.6),
                               SLIDE_W - Inches(2.8), Inches(3.2))
    band.adjustments[0] = 0.12
    shape_no_outline(band); fill_solid(band, PEACH)

    add_text(s, "Jr(t)  =  Jr,max  ·  ( 1 − e^(−λ · σ_t) )",
             left=Inches(1.4), top=Inches(3.8),
             width=SLIDE_W - Inches(2.8), height=Inches(1.8),
             size=80, bold=True, color=WHITE, font=TITLE_FONT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s,
             "σ_t  =  swarm diameter  =  mean over dimensions of  std(x) / box width",
             left=Inches(1.4), top=Inches(5.6),
             width=SLIDE_W - Inches(2.8), height=Inches(0.8),
             size=20, italic=True, color=WHITE, font=BODY_FONT,
             align=PP_ALIGN.CENTER)

    # behaviour cards
    cards = [
        ("01", BLUE,
         "Swarm spread out",
         "σ_t is large  →  Jr(t) climbs toward Jr,max\n"
         "We jump aggressively to keep exploring."),
        ("02", OLIVE,
         "Swarm converging",
         "σ_t shrinks  →  Jr(t) decays toward zero\n"
         "We stop spending FEs on OBL when it stops helping."),
        ("03", LAVEND,
         "Why it is new",
         "No reviewed paper in the OBL-PSO line\n"
         "ties Jr to a closed-form measurement of live\n"
         "swarm state.  This is our novel contribution."),
    ]
    add_three_cards(s, cards, top=Inches(7.4), height=Inches(3.5))

    s.notes_slide.notes_text_frame.text = (
        "Saqlain / Aena -- THE novel slide. Spend a beat on the two "
        "regimes. The 'Adaptive Jr in action' slide later will show this "
        "actually behaving as the formula predicts.")


# ---------- VA-5 + VA-6 ---------------------------------------------------

def slide_va_memory(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_background(s)
    add_content_header(s, "03", "VA-5 + VA-6",
                        "Memory and repulsion",
                        section_color=BLUE)
    add_subtitle(s, "Fix premature commitment at its root.")
    cards = [
        (BLUE,
         "VA-5   ·   pbest opposition",
         "Each particle keeps a stagnation counter s_i.\n"
         "When s_i ≥ t_stag = 5,\n"
         "    p̂_i = GOBL(p_i)\n"
         "Replace p_i with p̂_i if it is better.\n\n"
         "Attacks the root cause of premature\n"
         "convergence — a corrupted personal-best\n"
         "archive — rather than its symptoms."),
        (OLIVE,
         "VA-6   ·   Rebel repulsion",
         "Velocity update gains a repulsion term:\n"
         "    v_i ← χ [ … −α (p_worst − x_i) ]\n"
         "α = 0.05 (small but consistent).\n\n"
         "Particles are pushed AWAY from the swarm's\n"
         "currently worst pbest.\n"
         "Prevents the swarm from locking onto poor\n"
         "candidate basins."),
    ]
    add_two_cards(s, cards, top=Inches(3.6), height=Inches(6.6))

    s.notes_slide.notes_text_frame.text = (
        "Saqlain -- VA-5 is the most under-used trick in the literature; "
        "VA-6 is a tiny coefficient with a real effect at D=30.")


# ---------- VA-7 ----------------------------------------------------------

def slide_va_honest(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_background(s)
    add_content_header(s, "03", "VA-7",
                        "Honest reporting",
                        section_color=OLIVE)
    add_subtitle(s, "Function evaluations and statistical significance.")
    cards = [
        ("01", BLUE,
         "Function evaluations (FEs)",
         "Every OBL step charges an extra evaluation\n"
         "per particle. Reported in FEs, not\n"
         "iterations, so the cost is honest."),
        ("02", OLIVE,
         "30 independent runs",
         "Per (algorithm × benchmark × dim) cell.\n"
         "Seeds are fixed for exact reproducibility.\n"
         "Total: 1080 individual optimization runs."),
        ("03", LAVEND,
         "Wilcoxon signed-rank test",
         "One-sided alternative\n"
         "    AO-PSO < baseline.\n"
         "Reject H₀ at p < 0.05  →  significant win."),
    ]
    add_three_cards(s, cards, top=Inches(3.6), height=Inches(5.7))

    add_banner(s,
               "In our Phase-4 review, most papers did none of the above.",
               "Claims that look big can disappear under proper statistical testing.",
               top=Inches(9.6), color=DEEP)

    s.notes_slide.notes_text_frame.text = (
        "Inshal -- explain why we changed the x-axis to FEs and why "
        "Wilcoxon (not t-test) is correct for non-Gaussian fitness data.")


# ---------- Architecture --------------------------------------------------

def add_image_centered(slide, image_path, *, top, max_width=None,
                        max_height=None, left=None):
    from PIL import Image as PilImage
    with PilImage.open(image_path) as im:
        iw, ih = im.size
    if max_width is None:
        max_width = SLIDE_W - 2 * Inches(1.4)
    if max_height is None:
        max_height = SLIDE_H - top - Inches(0.6)
    img_w_emu = iw * 9525
    img_h_emu = ih * 9525
    scale = min(max_width / img_w_emu, max_height / img_h_emu)
    w = int(img_w_emu * scale)
    h = int(img_h_emu * scale)
    if left is None:
        left = (SLIDE_W - w) // 2
    slide.shapes.add_picture(image_path, left, top, width=w, height=h)


def slide_architecture(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_background(s)
    add_content_header(s, "04", "ARCHITECTURE",
                        "Where the seven enhancements live",
                        section_color=PEACH)
    add_subtitle(s,
                 "Stage 1 acts once; Stage 2 fires every generation; VA-7 wraps the experiment.")
    add_image_centered(s, fig("architecture.png"),
                        top=Inches(3.5))
    s.notes_slide.notes_text_frame.text = (
        "Eman -- walk through Stage 1 and Stage 2. Highlight VA-4 (the "
        "red novel block) in the middle of Stage 2.")


def slide_flowchart(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_background(s)
    add_content_header(s, "04", "PER-GENERATION FLOW",
                        "How AO-PSO steps through one iteration",
                        section_color=PEACH)
    add_subtitle(s, "Diamonds are decisions; rounded rectangles are processes.")

    add_image_centered(s, fig("flowchart.png"),
                        top=Inches(3.4),
                        max_width=Inches(7.5),
                        max_height=Inches(7.4),
                        left=Inches(1.4))

    # right side notes
    add_multiline(s, [
        ("Reading the flowchart",
         {"size": 28, "bold": True, "color": DEEP, "font": TITLE_FONT,
          "space_after": 16}),
        ("If U(0,1) < Jr(t)   →   generation jump (VA-3)",
         {"size": 19, "color": INK, "space_after": 8}),
        ("Any particle with s_i ≥ t_stag   →   pbest opposition (VA-5)",
         {"size": 19, "color": INK, "space_after": 8}),
        ("FE budget not exhausted   →   back to velocity update",
         {"size": 19, "color": INK, "space_after": 24}),
        ("No arrow crosses any node.",
         {"size": 22, "bold": True, "color": OLIVE, "font": TITLE_FONT,
          "space_after": 4}),
        ("Strictly vertical flow keeps the figure publication-clean.",
         {"size": 17, "italic": True, "color": GRAY}),
    ], left=Inches(10.0), top=Inches(4.2),
        width=Inches(8.5), height=Inches(6.0))


# ---------- experimental setup --------------------------------------------

def slide_setup(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_background(s)
    add_content_header(s, "05", "EXPERIMENTAL SETUP",
                        "What we ran, on what, and how many times",
                        section_color=SKY)

    # big numbers row -- the "business snapshot" style
    stats = [
        ("3",    "Algorithms",   "PSO, O-PSO, AO-PSO",      BLUE),
        ("6",    "Benchmarks",   "Sphere → Schwefel",       OLIVE),
        ("2",    "Dimensions",   "D = 10 and D = 30",       LAVEND),
        ("30",   "Seeds each",   "independent runs",        PEACH),
        ("1080", "Total runs",   "fully reproducible",      SKY),
    ]
    add_stat_row(s, stats, top=Inches(3.8))

    # bottom block: protocol details
    bx = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(1.4), Inches(7.7),
                             SLIDE_W - Inches(2.8), Inches(3.0))
    bx.adjustments[0] = 0.05
    shape_no_outline(bx); fill_solid(bx, WHITE)
    stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(1.4), Inches(7.7),
                                 Inches(0.25), Inches(3.0))
    shape_no_outline(stripe); fill_solid(stripe, DEEP)

    add_text(s, "Protocol",
             left=Inches(1.9), top=Inches(7.9),
             width=Inches(17), height=Inches(0.6),
             size=24, bold=True, color=DEEP, font=TITLE_FONT)
    add_multiline(s, [
        ("•  N = 20    c1 = c2 = 2.05    χ ≈ 0.7298    v_max = 0.5 (b−a)",
         {"size": 19, "color": INK, "space_after": 4}),
        ("•  FE budget   20 000  at  D = 10        60 000  at  D = 30",
         {"size": 19, "color": INK, "space_after": 4}),
        ("•  AO-PSO extras:  Jr,max = 0.30    λ = 25    t_stag = 5    α = 0.05",
         {"size": 19, "color": INK, "space_after": 4}),
        ("•  Significance: one-sided Wilcoxon signed-rank, α = 0.05",
         {"size": 19, "color": INK, "space_after": 4}),
    ], left=Inches(1.9), top=Inches(8.6),
        width=SLIDE_W - Inches(3.8), height=Inches(2.0))

    s.notes_slide.notes_text_frame.text = (
        "Inshal -- emphasise that we matched O-PSO's exact benchmark set "
        "for direct comparability, then added two harder functions.")


def add_stat_row_compact(s, stats, *, top):
    """Compact horizontal stat row -- ~3 inches tall.  Fits inside an
    11.25-inch slide even when the upper half is full of cards."""
    n = len(stats)
    gap = Inches(0.3)
    total_w = SLIDE_W - 2 * Inches(1.4)
    cell_w = (total_w - (n - 1) * gap) // n
    left0 = Inches(1.4)
    circle_d = Inches(1.7)
    for i, (number, label, sub, color) in enumerate(stats):
        left = left0 + i * (cell_w + gap)
        circ_left = left + (cell_w - circle_d) // 2
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                   circ_left, top, circle_d, circle_d)
        shape_no_outline(circ); fill_solid(circ, color)
        add_text(s, number,
                 left=left, top=top + Inches(0.15),
                 width=cell_w, height=circle_d - Inches(0.3),
                 size=36, bold=True, color=WHITE, font=TITLE_FONT,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, label,
                 left=left, top=top + circle_d + Inches(0.1),
                 width=cell_w, height=Inches(0.4),
                 size=16, bold=True, color=DEEP, font=TITLE_FONT,
                 align=PP_ALIGN.CENTER)
        add_text(s, sub,
                 left=left, top=top + circle_d + Inches(0.6),
                 width=cell_w, height=Inches(0.4),
                 size=13, italic=True, color=GRAY, font=BODY_FONT,
                 align=PP_ALIGN.CENTER)


def add_stat_row(s, stats, *, top, h=Inches(3.5)):
    n = len(stats)
    gap = Inches(0.3)
    total_w = SLIDE_W - 2 * Inches(1.4)
    cell_w = (total_w - (n - 1) * gap) // n
    left0 = Inches(1.4)
    for i, (number, label, sub, color) in enumerate(stats):
        left = left0 + i * (cell_w + gap)
        # colored circle bg
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                   left + (cell_w - Inches(2.4))//2,
                                   top, Inches(2.4), Inches(2.4))
        shape_no_outline(circ); fill_solid(circ, color)
        # big number
        add_text(s, number,
                 left=left, top=top + Inches(0.4),
                 width=cell_w, height=Inches(1.7),
                 size=60, bold=True, color=WHITE, font=TITLE_FONT,
                 align=PP_ALIGN.CENTER)
        # label below
        add_text(s, label,
                 left=left, top=top + Inches(2.6),
                 width=cell_w, height=Inches(0.4),
                 size=18, bold=True, color=DEEP, font=TITLE_FONT,
                 align=PP_ALIGN.CENTER)
        # sub
        add_text(s, sub,
                 left=left, top=top + Inches(3.0),
                 width=cell_w, height=Inches(0.4),
                 size=14, italic=True, color=GRAY, font=BODY_FONT,
                 align=PP_ALIGN.CENTER)


# ---------- results table -------------------------------------------------

def slide_results_table(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_background(s)
    add_content_header(s, "06", "RESULTS",
                        "Mean final fitness over 30 runs",
                        section_color=BLUE)
    add_subtitle(s, "Lower is better. Bold = AO-PSO dominates both baselines.")

    headers = ["Benchmark", "D", "Standard PSO", "O-PSO (focal)", "AO-PSO (ours)"]
    data = [
        ("Sphere",     10, "1.02 × 10⁻⁴⁰", "4.32 × 10⁻⁴¹", "5.62 × 10⁻³⁷",   False),
        ("Sphere",     30, "1.67 × 10⁺³",  "6.67 × 10⁺²",  "4.09 × 10⁻¹⁰⁶", True),
        ("Rosenbrock", 10, "6.24 × 10⁺³",  "3.06 × 10⁺³",  "4.12",            True),
        ("Rosenbrock", 30, "1.52 × 10⁺⁴",  "6.26 × 10⁺³",  "2.35 × 10⁺¹",    True),
        ("Rastrigin",  10, "1.06 × 10⁺¹",  "9.72",         "0.00",             True),
        ("Rastrigin",  30, "1.06 × 10⁺²",  "9.68 × 10⁺¹",  "0.00",             True),
        ("Griewank",   10, "8.67 × 10⁻²",  "1.09 × 10⁻¹",  "0.00",             True),
        ("Griewank",   30, "1.51 × 10⁺¹",  "6.38",         "0.00",             True),
        ("Ackley",     10, "2.86 × 10⁻¹",  "2.28 × 10⁻¹",  "4.44 × 10⁻¹⁶",    True),
        ("Ackley",     30, "4.47",         "4.46",         "4.44 × 10⁻¹⁶",    True),
        ("Schwefel",   10, "7.85 × 10⁺²",  "8.10 × 10⁺²",  "8.26 × 10⁺²",     False),
        ("Schwefel",   30, "4.09 × 10⁺³",  "3.86 × 10⁺³",  "5.57 × 10⁺³",     False),
    ]
    left = Inches(1.4)
    top = Inches(3.6)
    rows = len(data) + 1
    cols = 5
    tbl_shape = s.shapes.add_table(rows, cols, left, top,
                                    SLIDE_W - Inches(2.8), Inches(7.3))
    tbl = tbl_shape.table
    col_w = [Inches(3.0), Inches(1.0), Inches(4.2), Inches(4.2), Inches(4.8)]
    for i, w in enumerate(col_w):
        tbl.columns[i].width = w

    for c_idx, h in enumerate(headers):
        c = tbl.cell(0, c_idx); c.text = ""
        p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = h
        set_run(run, size=18, bold=True, color=WHITE, font=BODY_FONT)
        c.fill.solid(); c.fill.fore_color.rgb = DEEP

    for r_idx, row in enumerate(data, start=1):
        bench, dim, pso_v, opso_v, ao_v, bold_ao = row
        for c_idx, val in enumerate([bench, str(dim), pso_v, opso_v, ao_v]):
            cell = tbl.cell(r_idx, c_idx); cell.text = ""
            p = cell.text_frame.paragraphs[0]
            if c_idx == 0: p.alignment = PP_ALIGN.LEFT
            elif c_idx == 1: p.alignment = PP_ALIGN.CENTER
            else: p.alignment = PP_ALIGN.RIGHT
            r = p.add_run(); r.text = val
            ao_cell = c_idx == 4
            highlight = ao_cell and bold_ao
            set_run(r, size=15,
                    bold=highlight,
                    color=PEACH if highlight else INK,
                    font=BODY_FONT)
            cell.fill.solid()
            if r_idx % 2:
                cell.fill.fore_color.rgb = SKY
            else:
                cell.fill.fore_color.rgb = WHITE
            if ao_cell:
                # subtle peach background for the AO-PSO column
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xE8, 0xD8) if not highlight else RGBColor(0xFF, 0xD8, 0xC0)

    s.notes_slide.notes_text_frame.text = (
        "Inshal -- pause on Sphere D=30 (1e-106 vs 6.7e+02) and on "
        "Rastrigin and Griewank reaching exact zero. Acknowledge Schwefel.")


# ---------- Wilcoxon ------------------------------------------------------

def slide_wilcoxon(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_background(s)
    add_content_header(s, "06", "STATISTICAL SIGNIFICANCE",
                        "Wilcoxon signed-rank test",
                        section_color=BLUE)
    add_subtitle(s,
                 "p-value for the hypothesis  AO-PSO better than baseline.   Green ✓ = p < 0.05.")
    headers = ["Benchmark", "D", "vs.  Standard PSO", "vs.  O-PSO"]
    data = [
        ("Sphere",     10, "1.0",            "1.0"),
        ("Sphere",     30, "8.7 × 10⁻⁷ ✓",   "9.3 × 10⁻¹⁰ ✓"),
        ("Rosenbrock", 10, "0.627",           "0.110"),
        ("Rosenbrock", 30, "0.006 ✓",        "0.027 ✓"),
        ("Rastrigin",  10, "8.6 × 10⁻⁷ ✓",   "8.5 × 10⁻⁷ ✓"),
        ("Rastrigin",  30, "9.3 × 10⁻¹⁰ ✓",  "9.3 × 10⁻¹⁰ ✓"),
        ("Griewank",   10, "9.3 × 10⁻¹⁰ ✓",  "9.3 × 10⁻¹⁰ ✓"),
        ("Griewank",   30, "1.9 × 10⁻⁶ ✓",   "1.3 × 10⁻⁶ ✓"),
        ("Ackley",     10, "3.5 × 10⁻⁷ ✓",   "2.5 × 10⁻⁷ ✓"),
        ("Ackley",     30, "9.3 × 10⁻¹⁰ ✓",  "9.3 × 10⁻¹⁰ ✓"),
        ("Schwefel",   10, "0.86",            "0.63"),
        ("Schwefel",   30, "≈ 1.0",           "≈ 1.0"),
    ]
    left = Inches(1.6)
    top = Inches(3.6)
    rows = len(data) + 1
    tbl_shape = s.shapes.add_table(rows, 4, left, top,
                                    SLIDE_W - Inches(3.2), Inches(7.3))
    tbl = tbl_shape.table
    for i, w in enumerate([Inches(4.0), Inches(1.0), Inches(5.9), Inches(5.9)]):
        tbl.columns[i].width = w
    for c_idx, h in enumerate(headers):
        c = tbl.cell(0, c_idx); c.text = ""
        p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = h
        set_run(run, size=18, bold=True, color=WHITE, font=BODY_FONT)
        c.fill.solid(); c.fill.fore_color.rgb = DEEP
    for r_idx, row in enumerate(data, start=1):
        for c_idx, val in enumerate(row):
            cell = tbl.cell(r_idx, c_idx); cell.text = ""
            p = cell.text_frame.paragraphs[0]
            sig = "✓" in str(val)
            if c_idx == 0: p.alignment = PP_ALIGN.LEFT
            elif c_idx == 1: p.alignment = PP_ALIGN.CENTER
            else: p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = str(val)
            color = OLIVE if sig else GRAY
            set_run(r, size=16, bold=sig, color=color, font=BODY_FONT)
            cell.fill.solid()
            cell.fill.fore_color.rgb = SKY if r_idx % 2 else WHITE

    s.notes_slide.notes_text_frame.text = (
        "Inshal -- green check is significant. We pass on 5 of 6 benchmarks "
        "at D=30. The Schwefel row is the only loss; we deal with it later.")


# ---------- convergence plot ----------------------------------------------

def slide_convergence(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_background(s)
    add_content_header(s, "06", "CONVERGENCE",
                        "Median best fitness over 30 runs at D = 30",
                        section_color=BLUE)
    add_subtitle(s,
                 "AO-PSO (solid red) drops several orders of magnitude faster than either baseline.")
    add_image_centered(s, fig("convergence_D30.png"),
                        top=Inches(3.5))


def slide_boxplots(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_background(s)
    add_content_header(s, "06", "FINAL-FITNESS DISTRIBUTION",
                        "Box plots over 30 runs at D = 30",
                        section_color=BLUE)
    add_subtitle(s,
                 "AO-PSO has a lower median AND a dramatically tighter spread on multimodal functions.")
    add_image_centered(s, fig("boxplots_D30.png"), top=Inches(3.5))


def slide_diversity_jr(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_background(s)
    add_content_header(s, "06", "ADAPTIVE Jr IN ACTION",
                        "Rastrigin, D = 30, median of 30 runs",
                        section_color=PEACH)
    add_subtitle(s,
                 "Left: AO-PSO keeps the swarm diverse longer.   Right: Jr peaks early, then collapses to zero.")
    add_image_centered(s, fig("diversity_jr.png"), top=Inches(3.4))


# ---------- Schwefel honest limitation ------------------------------------

def slide_schwefel(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_background(s)
    add_content_header(s, "07", "HONEST LIMITATION",
                        "Where AO-PSO loses — and why",
                        section_color=OLIVE)
    add_subtitle(s, "On Schwefel at D = 30, AO-PSO ends at 5571 vs ~3900 for the baselines.")

    cards = [
        ("01", BLUE,
         "Schwefel's optimum",
         "x*  =  ( 420.97, ..., 420.97 )\n\n"
         "FAR from the centre of the search\n"
         "box  [−500, 500]^D."),
        ("02", OLIVE,
         "Why GOBL hurts",
         "Generalized OBL reflects across the box\n"
         "centre — straight away from the cluster\n"
         "of true optima at  ≈ 420."),
        ("03", LAVEND,
         "Why rebel hurts",
         "Worst pbest sits at the mirror local\n"
         "minimum.  Repulsion from it pushes\n"
         "particles toward the centre — wrong way."),
    ]
    add_three_cards(s, cards, top=Inches(3.6), height=Inches(5.0))

    add_banner(s,
               "This is a known but rarely-stated weakness of opposition-based methods.",
               "Future work:  centre-shifted opposite,   x̂ = 2 · x̄ − x.",
               top=Inches(9.0), color=DEEP)

    s.notes_slide.notes_text_frame.text = (
        "Aena -- own the loss. This honesty is what our literature review "
        "said the OBL-PSO field needs more of.")


# ---------- Conclusion / takeaways ----------------------------------------

def slide_conclusion(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_background(s)
    add_content_header(s, "08", "KEY TAKEAWAYS",
                        "What AO-PSO contributes",
                        section_color=LAVEND)

    cards = [
        ("01", BLUE,
         "Field-wide gaps closed",
         "OBL throughout the search.\n"
         "FE budget accounting.\n"
         "Wilcoxon signed-rank significance."),
        ("02", PEACH,
         "Novel contribution",
         "The adaptive jumping rate\n"
         "Jr(t) = Jr,max ( 1 − e^(−λσ_t) ).\n\n"
         "No prior OBL-PSO paper does this."),
        ("03", OLIVE,
         "Empirical result",
         "AO-PSO wins on 5 of 6 benchmarks at D=30\n"
         "Significant at  p < 10⁻⁶ on most\n"
         "Honest about Schwefel."),
    ]
    add_three_cards(s, cards, top=Inches(3.4), height=Inches(4.2))

    # bottom strip with headline stats -- compact, fits within slide height
    stats = [
        ("10¹⁰⁸",  "Sphere D=30 ratio",   "AO-PSO vs O-PSO",        BLUE),
        ("0.00",   "Rastrigin & Griewank","exact global optimum",   OLIVE),
        ("p<10⁻⁹","Wilcoxon",            "on most benchmarks",     LAVEND),
        ("1 / 6",  "Honest loss",        "Schwefel (off-centre)",  PEACH),
    ]
    add_stat_row_compact(s, stats, top=Inches(8.0))


# ---------- Thanks slide --------------------------------------------------

def slide_thanks(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_background(s)

    # decorative big shapes
    big = s.shapes.add_shape(MSO_SHAPE.OVAL,
                              -Inches(4.5), -Inches(4.5),
                              Inches(11), Inches(11))
    shape_no_outline(big); fill_solid(big, OLIVE)
    big2 = s.shapes.add_shape(MSO_SHAPE.OVAL,
                               SLIDE_W - Inches(6.5),
                               SLIDE_H - Inches(6.5),
                               Inches(10), Inches(10))
    shape_no_outline(big2); fill_solid(big2, BLUE)

    accent = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(15), Inches(2.0),
                                 Inches(1.0), Inches(1.0))
    shape_no_outline(accent); fill_solid(accent, PEACH)
    accent2 = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(3.5), Inches(8.5),
                                  Inches(0.7), Inches(0.7))
    shape_no_outline(accent2); fill_solid(accent2, LAVEND)

    # central card
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(3.5), Inches(3.2),
                               Inches(13.0), Inches(5.0))
    card.adjustments[0] = 0.06
    shape_no_outline(card); fill_solid(card, WHITE)

    add_text(s, "THANK YOU",
             left=Inches(3.5), top=Inches(3.7),
             width=Inches(13), height=Inches(2.5),
             size=140, bold=True, color=DEEP, font=TITLE_FONT,
             align=PP_ALIGN.CENTER)
    add_text(s, "Questions?",
             left=Inches(3.5), top=Inches(6.4),
             width=Inches(13), height=Inches(0.8),
             size=36, color=BLUE, font=TITLE_FONT,
             align=PP_ALIGN.CENTER)
    add_text(s, "AO-PSO  ·  Swarm Intelligence, Final Phase  ·  NUTECH",
             left=Inches(3.5), top=Inches(7.2),
             width=Inches(13), height=Inches(0.6),
             size=18, italic=True, color=GRAY, font=BODY_FONT,
             align=PP_ALIGN.CENTER)


# -------- driver -----------------------------------------------------------

def build() -> str:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_members(prs)
    slide_agenda(prs)

    slide_problem(prs)
    slide_related(prs)
    slide_focal(prs)

    slide_va_overview(prs)
    slide_va_init(prs)
    slide_va_gen_jumping(prs)
    slide_va_novel(prs)
    slide_va_memory(prs)
    slide_va_honest(prs)

    slide_architecture(prs)
    slide_flowchart(prs)

    slide_setup(prs)

    slide_results_table(prs)
    slide_wilcoxon(prs)
    slide_convergence(prs)
    slide_boxplots(prs)
    slide_diversity_jr(prs)

    slide_schwefel(prs)
    slide_conclusion(prs)
    slide_thanks(prs)

    prs.save(OUT_PATH)
    return OUT_PATH


if __name__ == "__main__":
    p = build()
    print(f"wrote {p}")
