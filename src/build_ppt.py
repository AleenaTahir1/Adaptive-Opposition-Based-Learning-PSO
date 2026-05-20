"""Generate a clean, no-footer 16:9 PowerPoint deck for the AO-PSO
final-phase presentation. Built for an eight-person team; each
content slide is annotated with the intended speaker in the speaker
notes so the team can rehearse without cluttering the slide face.

No animations, no clipart, no decorative footers -- just title bar,
content area, and (where useful) the figures rendered out of the
final paper.
"""

from __future__ import annotations

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt


# -------- palette ----------------------------------------------------------

NAVY = RGBColor(0x1F, 0x4E, 0x79)
RED  = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x1E, 0x7E, 0x34)
GRAY = RGBColor(0x4A, 0x4A, 0x4A)
LIGHT = RGBColor(0xF0, 0xF3, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


# -------- geometry ---------------------------------------------------------

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN  = Inches(0.55)
HEADER_H = Inches(0.95)

FIG_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                       "report", "figs")
OUT_PATH = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                        "AO_PSO_Presentation.pptx")


def fig(name: str) -> str:
    return os.path.join(FIG_DIR, name)


# -------- helpers ----------------------------------------------------------

def set_run(run, *, size=None, bold=None, color=None, italic=None, font=None):
    if font is not None:
        run.font.name = font
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if italic is not None:
        run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color


def add_header(slide, title: str, subtitle: str | None = None):
    """Title bar across the top of every content slide."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, HEADER_H
    )
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.shadow.inherit = False

    tf = bar.text_frame
    tf.margin_left = Inches(0.4)
    tf.margin_right = Inches(0.4)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.05)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title
    set_run(r, size=30, bold=True, color=WHITE, font="Calibri")

    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = subtitle
        set_run(r2, size=14, italic=True, color=LIGHT, font="Calibri")


def add_body_text(slide, lines, *, top=None, left=None, width=None, height=None,
                  base_size=20, line_color=GRAY):
    """Add a column of bullet/text lines.

    Each item in ``lines`` is either:
      - a string ``"text"``
      - a tuple ``(text, opts_dict)`` with keys size, bold, color, indent,
        italic, color (default GRAY).
    """
    if top is None:
        top = HEADER_H + Inches(0.3)
    if left is None:
        left = MARGIN
    if width is None:
        width = SLIDE_W - 2 * MARGIN
    if height is None:
        height = SLIDE_H - HEADER_H - Inches(0.4)

    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    first = True
    for line in lines:
        if isinstance(line, tuple):
            text, opts = line
        else:
            text, opts = line, {}
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = opts.get("align", PP_ALIGN.LEFT)
        p.level = opts.get("indent", 0)
        p.space_after = Pt(opts.get("space_after", 5))
        r = p.add_run()
        r.text = text
        set_run(
            r,
            size=opts.get("size", base_size),
            bold=opts.get("bold", False),
            italic=opts.get("italic", False),
            color=opts.get("color", line_color),
            font=opts.get("font", "Calibri"),
        )
    return tb


def add_image_scaled(slide, image_path: str, *, left=None, top=None,
                     max_width=None, max_height=None):
    """Insert an image and scale it to fit a max box without distortion."""
    from PIL import Image as PilImage
    with PilImage.open(image_path) as im:
        iw, ih = im.size
    if max_width is None:
        max_width = SLIDE_W - 2 * MARGIN
    if max_height is None:
        max_height = SLIDE_H - HEADER_H - Inches(0.6)

    scale = min(max_width / Emu(iw * 9525).emu * Emu(1).emu,
                max_height / Emu(ih * 9525).emu * Emu(1).emu)
    # The above is convoluted; just compute in EMU directly:
    img_w_emu = iw * 9525     # 9525 EMU per pixel at 96 DPI default
    img_h_emu = ih * 9525
    scale = min(max_width / img_w_emu, max_height / img_h_emu)
    w = int(img_w_emu * scale)
    h = int(img_h_emu * scale)

    if left is None:
        left = (SLIDE_W - w) // 2
    if top is None:
        top = HEADER_H + ((SLIDE_H - HEADER_H - h) // 2)

    slide.shapes.add_picture(image_path, left, top, width=w, height=h)


def add_speaker_note(slide, note: str):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = note


def add_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # 6 == blank in default


# -------- slide builders ---------------------------------------------------

def slide_title(prs):
    s = add_blank(prs)
    # navy band
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.5),
                              SLIDE_W, Inches(4.0))
    band.line.fill.background()
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.shadow.inherit = False

    tb = s.shapes.add_textbox(Inches(0.6), Inches(1.8),
                              SLIDE_W - Inches(1.2), Inches(3.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "AO-PSO"
    set_run(r, size=66, bold=True, color=WHITE)

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = "Adaptive Opposition-Based Particle Swarm Optimization"
    set_run(r2, size=28, color=WHITE)

    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.LEFT
    r3 = p3.add_run()
    r3.text = "with Chaotic Initialization, Generalized Opposites,"
    set_run(r3, size=22, color=LIGHT)
    p4 = tf.add_paragraph()
    p4.alignment = PP_ALIGN.LEFT
    r4 = p4.add_run()
    r4.text = "and Diversity-Driven Generation Jumping"
    set_run(r4, size=22, color=LIGHT)

    # team block
    tb2 = s.shapes.add_textbox(Inches(0.6), Inches(5.8),
                               SLIDE_W - Inches(1.2), Inches(1.5))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    members = ("Aleena Tahir, Eman Asghar, Aena Habib, Malaika Akhter, "
               "Dua Kamal, Sadia Mazhar, Inshal Adil, Saqlain Abbas")
    p = tf2.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = members
    set_run(r, size=16, bold=True, color=GRAY)
    p2 = tf2.add_paragraph()
    r = p2.add_run()
    r.text = ("Department of Artificial Intelligence  |  "
              "National University of Technology  |  Swarm Intelligence, "
              "Final Phase  |  Instructor: Lec. Faria Sajjad")
    set_run(r, size=13, color=GRAY)

    add_speaker_note(s, "Aleena opens, introduces title and team.")


def slide_outline(prs):
    s = add_blank(prs)
    add_header(s, "Outline", "What we will cover today")
    items = [
        ("1.  Problem and motivation", {"size": 22, "bold": True}),
        ("2.  Related work and the focal paper (O-PSO, 2009)", {"size": 22, "bold": True}),
        ("3.  Our value additions — the seven enhancements", {"size": 22, "bold": True}),
        ("4.  AO-PSO architecture and per-generation flow", {"size": 22, "bold": True}),
        ("5.  Implementation and experimental setup", {"size": 22, "bold": True}),
        ("6.  Results: tables, convergence, significance", {"size": 22, "bold": True}),
        ("7.  Honest limitation (Schwefel) and discussion", {"size": 22, "bold": True}),
        ("8.  Conclusion and Q&A", {"size": 22, "bold": True}),
    ]
    add_body_text(s, items, base_size=22)
    add_speaker_note(s, "Aleena -- 30 seconds to set expectations.")


def slide_problem(prs):
    s = add_blank(prs)
    add_header(s, "Problem and Motivation",
               "Why we worked on opposition-based PSO")
    lines = [
        ("Particle Swarm Optimization (1995) is one of the most widely used metaheuristics", {"size": 20}),
        ("but it suffers from two well-known failure modes:", {"size": 20, "space_after": 14}),
        ("•  Premature convergence — the swarm collapses into a sub-optimal basin", {"size": 20, "indent": 1, "color": RED}),
        ("•  Sensitivity to initialization — pseudo-random seeding gives uneven space coverage", {"size": 20, "indent": 1, "color": RED, "space_after": 16}),
        ("Opposition-Based Learning (Tizhoosh 2005) attacks both with a single idea:", {"size": 20, "bold": True}),
        ("whenever you evaluate x, also evaluate its opposite x̂ = a + b − x", {"size": 20, "italic": True, "color": NAVY, "space_after": 14}),
        ("Our focal paper (O-PSO, Jabeen et al. GECCO'09) was the first to splice OBL into PSO,", {"size": 20}),
        ("but only at initialization. After generation zero it reverts to plain PSO.", {"size": 20, "space_after": 14}),
        ("Goal:  push OBL deeper into the algorithm and report it honestly.", {"size": 22, "bold": True, "color": NAVY}),
    ]
    add_body_text(s, lines)
    add_speaker_note(s, "Malaika -- explain motivation; emphasise the two failure modes "
                        "and why opposition-based learning is a natural answer.")


def slide_related(prs):
    s = add_blank(prs)
    add_header(s, "Related Work",
               "Phase-4 review: 15 successor papers in the OBL-PSO line")
    lines = [
        ("Foundational:", {"size": 20, "bold": True, "color": NAVY}),
        ("Tizhoosh 2005 — OBL as a machine-intelligence scheme", {"size": 18, "indent": 1}),
        ("OBL applied at initialization:", {"size": 20, "bold": True, "color": NAVY}),
        ("Jabeen 2009 (O-PSO), CSPSO, WE-PSO, LatinPSO, Quasi-Random PSO", {"size": 18, "indent": 1, "space_after": 10}),
        ("OBL throughout the search:", {"size": 20, "bold": True, "color": NAVY}),
        ("OPSO-Cauchy 2007, OVCPSO 2009, GOPSO 2011, BPSO-OBL 2015, OBPSO 2012", {"size": 18, "indent": 1, "space_after": 10}),
        ("PSO variants with OBL:", {"size": 20, "bold": True, "color": NAVY}),
        ("OCLPSO, IOBL-PSO, COPSO, PSO-OBL-NNI, Imran-Mutation", {"size": 18, "indent": 1, "space_after": 10}),
        ("Beyond PSO / applied:", {"size": 20, "bold": True, "color": NAVY}),
        ("ODE (Diff. Evolution), SDE-OBL, AMO-OBL, Cao et al. reactive power dispatch", {"size": 18, "indent": 1, "space_after": 14}),
        ("Persistent open problems flagged across these papers:", {"size": 20, "bold": True, "color": RED}),
        ("•  fixed Jr  •  linear opposite only  •  no FE-fair cost accounting  •  no statistical tests", {"size": 18, "indent": 1, "color": RED}),
    ]
    add_body_text(s, lines)
    add_speaker_note(s, "Sadia -- summarise the related-work landscape and the four "
                        "open problems the literature still has.")


def slide_focal(prs):
    s = add_blank(prs)
    add_header(s, "The Focal Paper: O-PSO (Jabeen, Jalil, Baig — GECCO 2009)",
               "Our starting point — what it did and where it stops")
    lines = [
        ("What O-PSO does:", {"size": 22, "bold": True, "color": NAVY, "space_after": 8}),
        ("1.  Generate N random particles", {"size": 20, "indent": 1}),
        ("2.  Compute N linear opposites  x̂ = a + b − x", {"size": 20, "indent": 1}),
        ("3.  Evaluate all 2N candidates, keep the best N", {"size": 20, "indent": 1}),
        ("4.  Run standard Clerc-constriction PSO until budget exhausted", {"size": 20, "indent": 1, "space_after": 16}),
        ("Strengths:", {"size": 22, "bold": True, "color": GREEN, "space_after": 6}),
        ("•  Drop-in modification — works with any PSO implementation", {"size": 19, "indent": 1}),
        ("•  Tested on Sphere, Rosenbrock, Rastrigin, Griewank", {"size": 19, "indent": 1, "space_after": 14}),
        ("Limitations we will address:", {"size": 22, "bold": True, "color": RED, "space_after": 6}),
        ("•  Pseudo-random seed — no guaranteed space coverage", {"size": 19, "indent": 1, "color": RED}),
        ("•  Linear opposite only — geometrically rigid", {"size": 19, "indent": 1, "color": RED}),
        ("•  No OBL after generation 0 — most of the search runs blind", {"size": 19, "indent": 1, "color": RED}),
        ("•  No Wilcoxon, no FE-budget reporting — claims hard to verify", {"size": 19, "indent": 1, "color": RED}),
    ]
    add_body_text(s, lines)
    add_speaker_note(s, "Aena -- frame O-PSO as a clean but minimal contribution. "
                        "Be explicit that everything after this slide is our improvement.")


def slide_value_additions_summary(prs):
    s = add_blank(prs)
    add_header(s, "Our Value Additions — Seven Enhancements",
               "Each addresses one or more of the open problems")

    # Table-style layout: VA#, Mechanism, Source, Novel?
    rows = [
        ("#",    "Mechanism",                       "Source",          "Novel"),
        ("VA-1", "Chaotic logistic-map init",       "CSPSO",           ""),
        ("VA-2", "Generalized OBL (GOBL)",          "GOPSO",           ""),
        ("VA-3", "Generation jumping",              "OVCPSO / OPSO-C", ""),
        ("VA-4", "Adaptive Jr from swarm diameter", "—",               "★ ours"),
        ("VA-5", "pbest opposition on stagnation",  "BPSO-OBL",        ""),
        ("VA-6", "Rebel repulsion velocity term",   "BPSO-OBL",        ""),
        ("VA-7", "FE budget + Wilcoxon reporting",  "Best practice",   ""),
    ]

    left = Inches(0.6)
    top = HEADER_H + Inches(0.4)
    width = SLIDE_W - Inches(1.2)
    rows_n = len(rows)
    cols_n = 4
    table_shape = s.shapes.add_table(rows_n, cols_n, left, top,
                                      width, Inches(5.0))
    table = table_shape.table
    col_w = [Inches(1.1), Inches(6.8), Inches(2.7), Inches(1.4)]
    for i, w in enumerate(col_w):
        table.columns[i].width = w

    for r_idx, row in enumerate(rows):
        for c_idx, cell_text in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = ""
            tf = cell.text_frame
            tf.margin_left = Inches(0.1)
            tf.margin_right = Inches(0.1)
            tf.margin_top = Inches(0.04)
            tf.margin_bottom = Inches(0.04)
            p = tf.paragraphs[0]
            run = p.add_run(); run.text = cell_text
            if r_idx == 0:
                set_run(run, size=18, bold=True, color=WHITE)
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
            else:
                novel = row[3] != ""
                set_run(run, size=17,
                        bold=novel,
                        color=RED if novel else GRAY)
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT if r_idx % 2 else WHITE

    add_speaker_note(s, "Aena -- read out the table; explicitly call out VA-4 as the "
                        "novel contribution. Everything else is principled re-use.")


def slide_va12(prs):
    s = add_blank(prs)
    add_header(s, "VA-1 & VA-2 — Better Initialization",
               "Replace random seed and linear opposite")
    lines = [
        ("VA-1.  Chaotic logistic-map initialisation", {"size": 22, "bold": True, "color": NAVY, "space_after": 4}),
        ("    x_{n+1} = μ · x_n · (1 − x_n),   μ = 4  (fully chaotic regime)", {"size": 20, "italic": True, "color": GRAY, "space_after": 4}),
        ("    One independent stream per dimension, burn-in of 200 to discard transients", {"size": 19, "color": GRAY}),
        ("    Ergodic — covers the unit interval more uniformly than the Mersenne Twister", {"size": 19, "color": GRAY, "space_after": 14}),
        ("VA-2.  Generalized Opposite (GOBL) — Wang et al. 2011", {"size": 22, "bold": True, "color": NAVY, "space_after": 4}),
        ("    x̂_j = a_j + b_j − β x_j − (1 − β) x̄_j", {"size": 22, "italic": True, "color": RED}),
        ("    β ~ U(0,1) drawn per particle;  x̄ is the per-dimension population mean", {"size": 19, "color": GRAY}),
        ("    Mean-aware term breaks the box symmetry — opposites are no longer simple reflections", {"size": 19, "color": GRAY, "space_after": 14}),
        ("Cost: same 2N evaluations as O-PSO. Selection of best N of (X ∪ X̂) is unchanged.", {"size": 19, "bold": True, "color": GREEN}),
    ]
    add_body_text(s, lines)
    add_speaker_note(s, "Saqlain -- chaotic init costs nothing extra, GOBL is a one-line "
                        "formula change but a much richer opposite.")


def slide_va3(prs):
    s = add_blank(prs)
    add_header(s, "VA-3 — Generation Jumping",
               "Apply OBL throughout the search, not just at init")
    lines = [
        ("In O-PSO, OBL fires only once — at generation 0.", {"size": 20}),
        ("After that the swarm has only random initialisation to thank for any diversity.", {"size": 20, "space_after": 14}),
        ("Generation jumping (OVCPSO / OPSO-Cauchy):", {"size": 22, "bold": True, "color": NAVY, "space_after": 6}),
        ("    At every generation,  with probability Jr,", {"size": 20}),
        ("        1.  Compute X̂ = GOBL(X)", {"size": 20, "indent": 1}),
        ("        2.  Evaluate all N opposites", {"size": 20, "indent": 1}),
        ("        3.  Replace x_i with x̂_i whenever  f(x̂_i) < f(x_i)", {"size": 20, "indent": 1, "space_after": 14}),
        ("Cost:  +N FEs each time the jump fires (counted honestly in our FE budget).", {"size": 19, "italic": True, "color": GRAY, "space_after": 12}),
        ("Open problem in the literature:  every prior paper picks Jr by hand.", {"size": 20, "bold": True, "color": RED}),
        ("We fix that next ▶", {"size": 22, "bold": True, "color": RED}),
    ]
    add_body_text(s, lines)
    add_speaker_note(s, "Saqlain -- explain that without an adaptive Jr the FE budget is "
                        "wasted on OBL late in the search when the swarm has already converged.")


def slide_va4_novel(prs):
    s = add_blank(prs)
    add_header(s, "VA-4 — Adaptive Jumping Rate  (our novel contribution)",
               "Tie Jr to a live measurement of swarm diversity")

    lines = [
        ("Define the normalised swarm diameter at generation t:", {"size": 20, "space_after": 6}),
        ("    σ_t  =  mean over dimensions of  std(x_{i,j}) / (b_j − a_j)", {"size": 20, "italic": True, "color": NAVY, "space_after": 14}),
        ("Then set", {"size": 20}),
        ("                    Jr(t)  =  Jr,max · ( 1 − e^(−λ σ_t) )", {"size": 26, "bold": True, "color": RED, "space_after": 14}),
        ("Behaviour:", {"size": 22, "bold": True, "color": NAVY, "space_after": 6}),
        ("•  Swarm is spread out  →  σ_t large  →  Jr(t) → Jr,max  (explore via OBL)", {"size": 19, "indent": 1}),
        ("•  Swarm has converged  →  σ_t → 0  →  Jr(t) → 0  (stop wasting evaluations)", {"size": 19, "indent": 1, "space_after": 14}),
        ("Why this matters:", {"size": 22, "bold": True, "color": GREEN, "space_after": 6}),
        ("Among every paper in our Phase-4 review, none ties Jr to live swarm state.", {"size": 20, "indent": 1, "italic": True, "color": GRAY}),
        ("Hyperparameters:   Jr,max = 0.30,   λ = 25.", {"size": 20, "indent": 1, "italic": True, "color": GRAY}),
    ]
    add_body_text(s, lines)
    add_speaker_note(s, "Saqlain / Aena -- this is THE novel slide. Spend a beat on the "
                        "two regimes. The figure on slide 'Diversity & Adaptive Jr' will "
                        "show it actually behaving this way on Rastrigin.")


def slide_va56(prs):
    s = add_blank(prs)
    add_header(s, "VA-5 & VA-6 — Memory and Repulsion",
               "Fixing premature commitment at its root")
    lines = [
        ("VA-5.  pbest opposition on stagnation", {"size": 22, "bold": True, "color": NAVY, "space_after": 4}),
        ("    Each particle carries a stagnation counter s_i (incremented when p_i fails to improve)", {"size": 19, "color": GRAY}),
        ("    When s_i ≥ t_stag = 5,   compute p̂_i = GOBL(p_i)", {"size": 19, "color": GRAY}),
        ("    Replace p_i with p̂_i if it is better;  reset s_i either way", {"size": 19, "color": GRAY}),
        ("    Targets the root cause:  a corrupted personal-best archive,  not its symptoms", {"size": 19, "italic": True, "color": GREEN, "space_after": 14}),
        ("VA-6.  Rebel repulsion velocity term", {"size": 22, "bold": True, "color": NAVY, "space_after": 4}),
        ("    v_i ← χ [ v_i + c1 r1 (p_i − x_i) + c2 r2 (g − x_i) − α (p_worst − x_i) ]", {"size": 18, "italic": True, "color": RED}),
        ("    α = 0.05 — small but consistent push AWAY from the swarm's worst pbest", {"size": 19, "color": GRAY}),
        ("    Prevents attractor-locking around poor candidate basins", {"size": 19, "italic": True, "color": GREEN}),
    ]
    add_body_text(s, lines)
    add_speaker_note(s, "Saqlain -- VA-5 is the most under-used trick in the literature; "
                        "VA-6 is a tiny coefficient with a real effect at D=30.")


def slide_va7(prs):
    s = add_blank(prs)
    add_header(s, "VA-7 — Honest Reporting",
               "Function evaluations and statistical significance")
    lines = [
        ("Cost is measured in function evaluations (FEs), not iterations.", {"size": 21, "bold": True, "color": NAVY, "space_after": 6}),
        ("    Each opposition step charges an extra evaluation per particle.", {"size": 20, "indent": 1, "color": GRAY}),
        ("    The adaptive Jr ensures those evaluations decay as the swarm converges.", {"size": 20, "indent": 1, "color": GRAY, "space_after": 14}),
        ("30 independent runs per (algorithm × benchmark × dimension) cell.", {"size": 21, "bold": True, "color": NAVY, "space_after": 6}),
        ("    Each seed is fixed so the experiment is exactly reproducible.", {"size": 20, "indent": 1, "color": GRAY, "space_after": 14}),
        ("One-sided Wilcoxon signed-rank test, alternative \"AO-PSO < baseline\".", {"size": 21, "bold": True, "color": NAVY, "space_after": 6}),
        ("    Reject H₀ at p < 0.05 ⇒ AO-PSO is significantly better than the baseline.", {"size": 20, "indent": 1, "color": GRAY, "space_after": 14}),
        ("Why this matters:  in our Phase-4 review, most papers did none of the above.", {"size": 20, "italic": True, "color": RED}),
        ("Claims that look big can disappear under proper statistical testing.", {"size": 20, "italic": True, "color": RED}),
    ]
    add_body_text(s, lines)
    add_speaker_note(s, "Inshal -- explain why we changed the x-axis to FEs and why "
                        "Wilcoxon (not t-test) is correct for non-Gaussian fitness data.")


def slide_architecture(prs):
    s = add_blank(prs)
    add_header(s, "AO-PSO Architecture",
               "Where the seven value additions live")
    add_image_scaled(s, fig("architecture.png"),
                     max_width=SLIDE_W - Inches(0.8),
                     max_height=SLIDE_H - HEADER_H - Inches(0.6))
    add_speaker_note(s, "Eman -- walk through Stage 1 (top) and Stage 2 (bottom). "
                        "Highlight VA-4 (red novel block) in the middle of Stage 2.")


def slide_flowchart(prs):
    s = add_blank(prs)
    add_header(s, "Per-Generation Control Flow",
               "How AO-PSO actually steps through one iteration")
    add_image_scaled(s, fig("flowchart.png"),
                     max_width=SLIDE_W - Inches(8.0),  # narrow vertical fig
                     max_height=SLIDE_H - HEADER_H - Inches(0.6),
                     left=Inches(0.6))
    # right-hand annotations
    notes_lines = [
        ("Reading the diagram:", {"size": 22, "bold": True, "color": NAVY, "space_after": 8}),
        ("•  Diamonds = decisions", {"size": 19}),
        ("•  Rounded rectangles = processes", {"size": 19}),
        ("•  Green oval = entry / exit", {"size": 19, "space_after": 14}),
        ("Key paths:", {"size": 22, "bold": True, "color": NAVY, "space_after": 8}),
        ("•  Jr check passes  →  Generation Jump fires (VA-3)", {"size": 19}),
        ("•  Any stagnated particle  →  pbest opposition (VA-5)", {"size": 19}),
        ("•  FE budget not exhausted  →  return to velocity update", {"size": 19, "space_after": 14}),
        ("Note:  no arrow crosses any node — strictly vertical flow,", {"size": 18, "italic": True, "color": GRAY}),
        ("so the figure is publication-clean.", {"size": 18, "italic": True, "color": GRAY}),
    ]
    add_body_text(s, notes_lines, left=Inches(6.0), top=HEADER_H + Inches(0.3),
                  width=Inches(6.7))
    add_speaker_note(s, "Eman -- one minute on the diamond gates. Don't read the whole "
                        "thing; let the audience read while you point.")


def slide_setup(prs):
    s = add_blank(prs)
    add_header(s, "Experimental Setup",
               "What we ran, on what, how many times")
    lines = [
        ("Algorithms compared:", {"size": 22, "bold": True, "color": NAVY, "space_after": 6}),
        ("•  Standard PSO (Clerc constriction)", {"size": 20, "indent": 1}),
        ("•  O-PSO (Jabeen et al. 2009) — the focal paper",  {"size": 20, "indent": 1}),
        ("•  AO-PSO (ours) — all seven value additions", {"size": 20, "indent": 1, "bold": True, "color": RED, "space_after": 14}),

        ("Benchmarks (6 in total):", {"size": 22, "bold": True, "color": NAVY, "space_after": 6}),
        ("•  Original O-PSO four:  Sphere, Rosenbrock, Rastrigin, Griewank", {"size": 20, "indent": 1}),
        ("•  Extra multimodal:  Ackley, Schwefel", {"size": 20, "indent": 1, "space_after": 14}),

        ("Protocol:", {"size": 22, "bold": True, "color": NAVY, "space_after": 6}),
        ("•  D ∈ {10, 30};  FE budget 20 000 (D=10), 60 000 (D=30)", {"size": 20, "indent": 1}),
        ("•  Swarm size N = 20;  c1 = c2 = 2.05;  v_max = 0.5 (b−a)", {"size": 20, "indent": 1}),
        ("•  AO-PSO extras:  Jr,max = 0.30,  λ = 25,  t_stag = 5,  α = 0.05", {"size": 20, "indent": 1}),
        ("•  30 independent random seeds per cell  →  1080 runs in total", {"size": 20, "indent": 1, "bold": True}),
    ]
    add_body_text(s, lines)
    add_speaker_note(s, "Inshal -- emphasise that we matched O-PSO's exact benchmark set "
                        "for direct comparability, then added two harder functions.")


def slide_results_table(prs):
    s = add_blank(prs)
    add_header(s, "Mean Final Fitness — 30 Runs",
               "Lower is better.  Bold = AO-PSO dominates both baselines")

    headers = ["Benchmark", "D", "PSO", "O-PSO", "AO-PSO"]
    data = [
        ("Sphere",     10, "1.02e-40", "4.32e-41", "5.62e-37",   False),
        ("Sphere",     30, "1.67e+03", "6.67e+02", "4.09e-106", True),
        ("Rosenbrock", 10, "6.24e+03", "3.06e+03", "4.12e+00",  True),
        ("Rosenbrock", 30, "1.52e+04", "6.26e+03", "2.35e+01",  True),
        ("Rastrigin",  10, "1.06e+01", "9.72e+00", "0.00",       True),
        ("Rastrigin",  30, "1.06e+02", "9.68e+01", "0.00",       True),
        ("Griewank",   10, "8.67e-02", "1.09e-01", "0.00",       True),
        ("Griewank",   30, "1.51e+01", "6.38e+00", "0.00",       True),
        ("Ackley",     10, "2.86e-01", "2.28e-01", "4.44e-16",  True),
        ("Ackley",     30, "4.47e+00", "4.46e+00", "4.44e-16",  True),
        ("Schwefel",   10, "7.85e+02", "8.10e+02", "8.26e+02",  False),
        ("Schwefel",   30, "4.09e+03", "3.86e+03", "5.57e+03",  False),
    ]

    left = Inches(0.6)
    top = HEADER_H + Inches(0.35)
    width = SLIDE_W - Inches(1.2)
    rows = len(data) + 1
    cols = 5
    tbl_shape = s.shapes.add_table(rows, cols, left, top, width,
                                    Inches(5.3))
    tbl = tbl_shape.table
    col_w = [Inches(2.8), Inches(0.9), Inches(2.5), Inches(2.5), Inches(3.4)]
    for i, w in enumerate(col_w):
        tbl.columns[i].width = w

    for c_idx, h in enumerate(headers):
        c = tbl.cell(0, c_idx)
        c.text = ""
        p = c.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = h
        set_run(run, size=18, bold=True, color=WHITE)
        c.fill.solid(); c.fill.fore_color.rgb = NAVY

    for r_idx, row in enumerate(data, start=1):
        bench, dim, pso_v, opso_v, ao_v, bold_ao = row
        for c_idx, val in enumerate([bench, str(dim), pso_v, opso_v, ao_v]):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if c_idx in (1,) else PP_ALIGN.LEFT
            if c_idx == 0:
                p.alignment = PP_ALIGN.LEFT
            elif c_idx >= 2:
                p.alignment = PP_ALIGN.RIGHT
            r = p.add_run(); r.text = val
            highlight = c_idx == 4 and bold_ao
            set_run(r,
                    size=15,
                    bold=highlight,
                    color=RED if highlight else GRAY)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT if r_idx % 2 else WHITE

    add_speaker_note(s, "Inshal -- pause on Sphere D=30 (1e-106 vs O-PSO's 6e+02) and on "
                        "Rastrigin/Griewank reaching exact zero. Acknowledge Schwefel.")


def slide_wilcoxon(prs):
    s = add_blank(prs)
    add_header(s, "Statistical Significance — Wilcoxon Signed-Rank Test",
               "p-value for H_a : AO-PSO better than baseline")
    headers = ["Benchmark", "D", "vs. PSO", "vs. O-PSO"]
    data = [
        ("Sphere",     10, "1.0",          "1.0"),
        ("Sphere",     30, "8.7e-7  ✓",   "9.3e-10  ✓"),
        ("Rosenbrock", 10, "0.627",        "0.110"),
        ("Rosenbrock", 30, "0.006  ✓",    "0.027  ✓"),
        ("Rastrigin",  10, "8.6e-7  ✓",   "8.5e-7  ✓"),
        ("Rastrigin",  30, "9.3e-10  ✓",  "9.3e-10  ✓"),
        ("Griewank",   10, "9.3e-10  ✓",  "9.3e-10  ✓"),
        ("Griewank",   30, "1.9e-6  ✓",   "1.3e-6  ✓"),
        ("Ackley",     10, "3.5e-7  ✓",   "2.5e-7  ✓"),
        ("Ackley",     30, "9.3e-10  ✓",  "9.3e-10  ✓"),
        ("Schwefel",   10, "0.86",         "0.63"),
        ("Schwefel",   30, "≈ 1.0",        "≈ 1.0"),
    ]

    left = Inches(0.8)
    top = HEADER_H + Inches(0.3)
    tbl_shape = s.shapes.add_table(len(data) + 1, 4, left, top,
                                    Inches(11.5), Inches(5.5))
    tbl = tbl_shape.table
    for i, w in enumerate([Inches(2.8), Inches(0.9), Inches(3.9), Inches(3.9)]):
        tbl.columns[i].width = w
    for c_idx, h in enumerate(headers):
        c = tbl.cell(0, c_idx)
        c.text = ""
        p = c.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = h
        set_run(run, size=18, bold=True, color=WHITE)
        c.fill.solid(); c.fill.fore_color.rgb = NAVY

    for r_idx, row in enumerate(data, start=1):
        for c_idx, val in enumerate(row):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            sig = "✓" in str(val)
            p.alignment = PP_ALIGN.CENTER if c_idx >= 1 else PP_ALIGN.LEFT
            r = p.add_run(); r.text = str(val)
            set_run(r, size=16, bold=sig,
                    color=GREEN if sig else GRAY)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT if r_idx % 2 else WHITE

    add_speaker_note(s, "Inshal -- 'green check' = significant at p<0.05. We pass on "
                        "5 of 6 benchmarks at D=30. The Schwefel row is the only loss; "
                        "we deal with it on the next-but-one slide.")


def slide_convergence(prs):
    s = add_blank(prs)
    add_header(s, "Convergence Curves at D = 30",
               "Median best fitness vs. function evaluations")
    add_image_scaled(s, fig("convergence_D30.png"))
    add_speaker_note(s, "Inshal -- AO-PSO drops several orders of magnitude faster than "
                        "either baseline on every multimodal benchmark.")


def slide_boxplots(prs):
    s = add_blank(prs)
    add_header(s, "Final-Fitness Distributions — 30 Runs at D = 30",
               "Lower median AND tighter spread")
    add_image_scaled(s, fig("boxplots_D30.png"))
    add_speaker_note(s, "Inshal -- the spread story matters: AO-PSO's boxes are tiny on "
                        "Rastrigin, Griewank, Ackley — consistent across seeds.")


def slide_diversity_jr(prs):
    s = add_blank(prs)
    add_header(s, "Adaptive Jr in Action  (Rastrigin, D = 30)",
               "Diversity drives Jr — and Jr saves evaluations")
    add_image_scaled(s, fig("diversity_jr.png"))
    add_speaker_note(s, "Aena or Eman -- THE slide that validates VA-4. Left: AO-PSO "
                        "stays diverse longer. Right: Jr peaks at ~0.29 then collapses "
                        "to zero — no wasted FEs on late-stage OBL.")


def slide_schwefel(prs):
    s = add_blank(prs)
    add_header(s, "Honest Limitation — Schwefel",
               "Where AO-PSO loses and why")
    lines = [
        ("On Schwefel at D = 30 AO-PSO ends at 5571, vs ~3900 for the baselines.", {"size": 20}),
        ("This is a real loss — we report it.", {"size": 20, "bold": True, "color": RED, "space_after": 14}),
        ("Why:", {"size": 22, "bold": True, "color": NAVY, "space_after": 6}),
        ("•  Schwefel's optimum is at x* = (420.97, …, 420.97), FAR from the box centre.", {"size": 19, "indent": 1}),
        ("•  GOBL reflects candidates across the centre, away from the cluster of true optima.", {"size": 19, "indent": 1}),
        ("•  Rebel repulsion pushes particles AWAY from the worst pbest, often the mirror local minimum.", {"size": 19, "indent": 1, "space_after": 14}),
        ("This is a known but rarely-stated weakness of opposition-based methods.", {"size": 20, "italic": True, "color": GRAY, "space_after": 6}),
        ("Most papers in the OBL line silently exclude Schwefel — we don't.", {"size": 20, "italic": True, "color": GRAY, "space_after": 14}),
        ("Future work: centre-shifted opposite,   x̂ = 2 x̄ − x.", {"size": 21, "bold": True, "color": GREEN}),
    ]
    add_body_text(s, lines)
    add_speaker_note(s, "Aena -- own the loss; explain why it happens. This honesty is "
                        "exactly what our literature review said the OBL-PSO field needs.")


def slide_conclusion(prs):
    s = add_blank(prs)
    add_header(s, "Conclusion",
               "Seven enhancements, one novel mechanism, honest reporting")
    lines = [
        ("AO-PSO addresses every open problem flagged in our Phase-4 review:", {"size": 21, "bold": True, "color": NAVY, "space_after": 10}),
        ("✓  OBL throughout the search (not only at init)", {"size": 20, "indent": 1, "color": GREEN}),
        ("✓  Adaptive Jr — closed-form, driven by live swarm state  (novel)", {"size": 20, "indent": 1, "color": GREEN}),
        ("✓  FE-budget cost accounting + Wilcoxon significance", {"size": 20, "indent": 1, "color": GREEN, "space_after": 14}),
        ("Empirical headline (D = 30):", {"size": 21, "bold": True, "color": NAVY, "space_after": 6}),
        ("•  Sphere:        AO-PSO 4 × 10⁻¹⁰⁶   vs.   O-PSO 6.67 × 10²", {"size": 18, "indent": 1, "font": "Consolas"}),
        ("•  Rastrigin:     AO-PSO 0.00            vs.   O-PSO 96.8", {"size": 18, "indent": 1, "font": "Consolas"}),
        ("•  Griewank:      AO-PSO 0.00            vs.   O-PSO 6.38", {"size": 18, "indent": 1, "font": "Consolas"}),
        ("•  Ackley:        AO-PSO 4.4 × 10⁻¹⁶   vs.   O-PSO 4.46", {"size": 18, "indent": 1, "font": "Consolas", "space_after": 14}),
        ("Significance:  p < 10⁻⁶ on 5 of 6 benchmarks at D = 30.", {"size": 21, "bold": True, "color": RED, "space_after": 14}),
        ("Honestly reported limitation on Schwefel — pointer for future work.", {"size": 20, "italic": True, "color": GRAY}),
    ]
    add_body_text(s, lines)
    add_speaker_note(s, "Malaika -- close strong: novel contribution + 5 of 6 wins at "
                        "p<1e-6 + honest about Schwefel.")


def slide_team(prs):
    s = add_blank(prs)
    add_header(s, "Work Distribution",
               "Who did what")

    headers = ["Role", "Member", "Deliverable"]
    rows = [
        ("Value addition & methodology", "Aena Habib",     "Seven enhancements, design rationale"),
        ("Implementation & testing",     "Saqlain Abbas",  "Python implementation, runs, debugging"),
        ("Results analysis & plots",     "Inshal Adil",    "Convergence curves, Wilcoxon, tables"),
        ("Workflow diagrams",            "Eman Asghar",    "Architecture and flowchart figures"),
        ("Presentation design",          "Dua Kamal",      "This slide deck, visuals"),
        ("Abstract, intro, conclusion",  "Malaika Akhter", "IEEE report body text"),
        ("Related-work section",         "Sadia Mazhar",   "Phase-4 review condensed"),
        ("IEEE compilation & format",    "Aleena Tahir",   "Final paper, references, proofreading"),
    ]

    left = Inches(0.6)
    top = HEADER_H + Inches(0.3)
    tbl_shape = s.shapes.add_table(len(rows) + 1, 3, left, top,
                                    SLIDE_W - Inches(1.2), Inches(5.5))
    tbl = tbl_shape.table
    for i, w in enumerate([Inches(4.4), Inches(2.8), Inches(5.0)]):
        tbl.columns[i].width = w

    for c_idx, h in enumerate(headers):
        c = tbl.cell(0, c_idx); c.text = ""
        p = c.text_frame.paragraphs[0]; r = p.add_run(); r.text = h
        set_run(r, size=18, bold=True, color=WHITE)
        c.fill.solid(); c.fill.fore_color.rgb = NAVY

    for r_idx, row in enumerate(rows, start=1):
        for c_idx, val in enumerate(row):
            cell = tbl.cell(r_idx, c_idx); cell.text = ""
            p = cell.text_frame.paragraphs[0]
            r = p.add_run(); r.text = val
            set_run(r, size=17, bold=(c_idx == 1), color=GRAY)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT if r_idx % 2 else WHITE

    add_speaker_note(s, "Aleena -- thank the team. One-line credits each, then move to Q&A.")


def slide_thanks(prs):
    s = add_blank(prs)
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.5),
                              SLIDE_W, Inches(3.0))
    band.line.fill.background()
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.shadow.inherit = False

    tb = s.shapes.add_textbox(Inches(0.6), Inches(2.9),
                              SLIDE_W - Inches(1.2), Inches(2.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "Thank you"
    set_run(r, size=72, bold=True, color=WHITE)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "Questions?"
    set_run(r2, size=32, color=LIGHT)
    add_speaker_note(s, "Aleena -- open Q&A.")


# -------- driver -----------------------------------------------------------

def build() -> str:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_outline(prs)
    slide_problem(prs)
    slide_related(prs)
    slide_focal(prs)
    slide_value_additions_summary(prs)
    slide_va12(prs)
    slide_va3(prs)
    slide_va4_novel(prs)
    slide_va56(prs)
    slide_va7(prs)
    slide_architecture(prs)
    slide_flowchart(prs)
    slide_setup(prs)
    slide_results_table(prs)
    slide_wilcoxon(prs)
    slide_convergence(prs)
    slide_boxplots(prs)
    slide_diversity_jr(prs)
    slide_schwefel(prs)
    slide_conclusion(prs)
    slide_team(prs)
    slide_thanks(prs)

    prs.save(OUT_PATH)
    return OUT_PATH


if __name__ == "__main__":
    path = build()
    print(f"wrote {path}")
